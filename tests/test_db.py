"""
WEP AI — Test suite v14.7

Tests de regresión para fixes aplicados en v14.1 a v14.7.
Correr con:
    cd wepai/
    python3 tests/test_db.py

Cubre:
 1. init_db sin errores
 2. Registro y login con bcrypt
 3. Migración SHA-256 → bcrypt en login
 4. Plan free: límite de 5 docs/mes
 5. Race condition resuelta (20 threads paralelos)
 6. save_card_last4 con validación
 7. get_plan_limits para free/pro/biz
 8. CRUD de conversations / messages / documents
 9. Dispatcher template_id (Excel, v14.2) — precedencia y fallback
10. Word controller preserva datos del usuario (v14.3) — anti data-loss
11. Cobertura bilingüe ES/EN (v14.4) — Word/Excel/PowerPoint
12. Production-ready (v14.5) — fechas US, moneda USD, Stripe scaffolding
13. v14.6 — ui_lang fallback, email module, trial 14d, billing helpers
14. v14.7 — email integration, Stripe metadata, cancel direct, status banner

Adicionalmente:
- tests/test_mac_e2e.py — E2E real con macOS + Office (correr en Mac)
"""

import os
import sys
import tempfile
import threading
import hashlib
import sqlite3
import importlib.util

# Ruta al módulo db
HERE = os.path.dirname(os.path.abspath(__file__))
DB_PATH = os.path.join(HERE, "..", "storage", "db.py")


def _load_db_module(tmp_db_path):
    """Carga db.py apuntando a una BD temporal aislada."""
    spec = importlib.util.spec_from_file_location("db", DB_PATH)
    db = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(db)
    db.DB_PATH = tmp_db_path
    return db


def main():
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".db")
    tmp.close()
    db = _load_db_module(tmp.name)

    try:
        run_tests(db, tmp.name)
    finally:
        os.unlink(tmp.name)


def run_tests(db, tmp_path):
    print("== TEST 1: init_db ==")
    db.init_db()
    print("  init_db OK")

    print("\n== TEST 2: register + login con bcrypt ==")
    uid, err = db.register_user("Juan", "Perez", "juan@test.com", "MiPassword123!")
    assert err is None, f"register failed: {err}"
    assert uid is not None
    print(f"  registered uid={uid}")

    user = db.login_user("juan@test.com", "MiPassword123!")
    assert user is not None and user["id"] == uid
    print(f"  login OK")

    bad = db.login_user("juan@test.com", "WrongPassword")
    assert bad is None
    print("  login rechaza password incorrecto")

    print("\n== TEST 3: migración SHA-256 → bcrypt en login ==")
    legacy_hash = hashlib.sha256("LegacyPass".encode()).hexdigest()
    con = sqlite3.connect(tmp_path)
    con.execute(
        "INSERT INTO users (name, last_name, email, password_hash, plan, month_reset, created_at) "
        "VALUES (?,?,?,?,?,?,?)",
        ("Old", "User", "old@test.com", legacy_hash, "free", "2026-05", "2026-05-21"),
    )
    con.commit(); con.close()

    user = db.login_user("old@test.com", "LegacyPass")
    assert user is not None, "legacy login failed"

    con = sqlite3.connect(tmp_path)
    hash_now = con.execute(
        "SELECT password_hash FROM users WHERE email=?", ("old@test.com",)
    ).fetchone()[0]
    con.close()
    assert hash_now.startswith(("$2a$", "$2b$", "$2y$")), f"hash no migrado: {hash_now[:10]}"
    print(f"  hash auto-upgraded a bcrypt ({hash_now[:7]}...)")

    print("\n== TEST 4: plan free — 5 docs permitidos, 6° rechazado ==")
    for i in range(5):
        ok, reason = db.check_and_consume_doc(uid)
        assert ok, f"doc {i+1} debía aceptarse, got {reason}"
    print("  5 docs consumidos correctamente")

    ok, reason = db.check_and_consume_doc(uid)
    assert not ok and reason == "limit_reached"
    print("  doc 6 rechazado con limit_reached")

    print("\n== TEST 5: race condition bajo 20 threads ==")
    con = sqlite3.connect(tmp_path)
    con.execute("UPDATE users SET docs_used_this_month=0 WHERE id=?", (uid,))
    con.commit(); con.close()

    results = []
    lock = threading.Lock()

    def worker():
        ok, reason = db.check_and_consume_doc(uid)
        with lock:
            results.append((ok, reason))

    threads = [threading.Thread(target=worker) for _ in range(20)]
    for t in threads:
        t.start()
    for t in threads:
        t.join()

    successes = sum(1 for ok, _ in results if ok)
    rejections = sum(1 for ok, _ in results if not ok)
    print(f"  20 threads → {successes} aceptados, {rejections} rechazados")
    assert successes == 5, f"race condition: {successes} success (expected 5)"
    print("  race condition resuelta")

    print("\n== TEST 6: save_card_last4 con validación ==")
    db.save_card_last4(uid, "1234")
    con = sqlite3.connect(tmp_path)
    last4 = con.execute("SELECT card_last4 FROM users WHERE id=?", (uid,)).fetchone()[0]
    con.close()
    assert last4 == "1234"
    print(f"  card_last4 saved: {last4}")

    db.save_card_last4(uid, "abcd")  # debe rechazar
    con = sqlite3.connect(tmp_path)
    last4 = con.execute("SELECT card_last4 FROM users WHERE id=?", (uid,)).fetchone()[0]
    con.close()
    assert last4 == "1234"
    print("  no-digit values rechazados")

    print("\n== TEST 7: get_plan_limits ==")
    assert db.get_plan_limits("free")["docs_month"] == 5
    assert db.get_plan_limits("pro")["docs_month"] == -1
    assert db.get_plan_limits("biz")["docs_month"] == -1
    print("  plan limits OK")

    print("\n== TEST 8: conversations / messages / documents ==")
    cid = db.new_conversation(uid, "Test convo")
    db.save_message(cid, "user", "Hola")
    db.save_message(cid, "assistant", "Hola, en qué te ayudo?")
    msgs = db.get_messages(cid)
    assert len(msgs) == 2 and msgs[0]["role"] == "user"
    db.save_document(cid, uid, "Mi doc", "word", "/tmp/doc.docx")
    docs = db.get_documents(cid)
    assert len(docs) == 1
    print(f"  conversation cid={cid}, {len(msgs)} mensajes, {len(docs)} docs")

    print("\n== TEST 9: dispatcher template_id (Excel, v14.2) ==")
    # Importar el controller (necesita openpyxl)
    try:
        sys.path.insert(0, os.path.join(HERE, ".."))
        from office import excel_controller as ec
        from openpyxl import load_workbook
        import tempfile as _tf
        # Aislamos OUTPUT_DIR para no contaminar
        tmpdir = _tf.mkdtemp()
        ec.OUTPUT_DIR = tmpdir

        # template_id explícito gana sobre keyword
        fp, _ = ec.generate_excel({
            "tipo": "excel", "template_id": "dashboard",
            "titulo": "T", "instrucciones": "inventario stock",
            "datos": {"detalles_especificos": "", "secciones": []},
        })
        wb = load_workbook(fp); sheets = wb.sheetnames
        assert any("Dashboard" in s or "KPI" in s.upper() for s in sheets), sheets
        print(f"  template_id=dashboard ganó sobre 'inventario' keyword (sheets: {sheets})")
        os.unlink(fp)

        # fallback a keyword cuando no hay template_id
        fp, _ = ec.generate_excel({
            "tipo": "excel",
            "titulo": "T", "instrucciones": "inventario de productos",
            "datos": {"detalles_especificos": "", "secciones": []},
        })
        wb = load_workbook(fp); sheets = wb.sheetnames
        assert "Inventario" in sheets, sheets
        print(f"  sin template_id → fallback a keyword (sheets: {sheets})")
        os.unlink(fp)

        # template_id inválido → fallback
        fp, _ = ec.generate_excel({
            "tipo": "excel", "template_id": "no_existe",
            "titulo": "T", "instrucciones": "inventario stock",
            "datos": {"detalles_especificos": "", "secciones": []},
        })
        wb = load_workbook(fp); sheets = wb.sheetnames
        assert "Inventario" in sheets, sheets
        print(f"  template_id inválido → fallback a keyword OK")
        os.unlink(fp)
    except ImportError as e:
        print(f"  SKIP (openpyxl/excel_controller no disponible: {e})")

    print("\n== TEST 10: word_controller data preservation (v14.3) ==")
    try:
        from office import word_controller as wc
        from docx import Document
        import tempfile as _tf
        tmpdir = _tf.mkdtemp()
        wc.OUTPUT_DIR = tmpdir

        gen_data = {
            "tipo": "word", "template_id": "commercial_lease",
            "titulo": "Test_Lease",
            "datos": {
                "pais": "CO",
                "parte_a_nombre": "Juan Pérez Pérez",
                "parte_a_id": "CC 12345678",
                "parte_b_nombre": "Logística XYZ S.A.S.",
                "parte_b_id": "NIT 900123456-7",
                "objeto": "Bodega Carrera 50",
                "monto": "8.500.000",
                "duracion": "24 meses",
                "secciones": [],
            },
        }
        fp = wc.generate_word(gen_data)
        d = Document(fp)
        text = "\n".join(p.text for p in d.paragraphs)
        for t in d.tables:
            for row in t.rows:
                for cell in row.cells:
                    text += "\n" + cell.text

        critical_data = [
            "Juan Pérez Pérez", "Logística XYZ", "CC 12345678",
            "NIT 900123456-7", "Bodega Carrera 50", "8.500.000", "24 meses",
        ]
        missing = [d for d in critical_data if d not in text]
        assert not missing, f"data loss: {missing}"
        print(f"  Los 7 datos críticos del usuario aparecen en el doc")

        # Test 2: el _user_data_block aparece como red de seguridad
        assert "DATOS CLAVE DEL DOCUMENTO" in text, "red de seguridad no insertada"
        print(f"  Sección 'Datos clave del documento' insertada como red de seguridad")
        os.unlink(fp)
    except ImportError as e:
        print(f"  SKIP (python-docx/word_controller no disponible: {e})")

    # ── TEST 11 — v14.4: Bilingüe ES/EN en Word, Excel y PowerPoint ───────────
    print("\n[11] v14.4: cobertura bilingüe ES/EN")
    try:
        import tempfile as _tf
        from office import word_controller as wc
        from office import excel_controller as ec
        from office import ppt_controller as pc
        from docx import Document
        from openpyxl import load_workbook
        from pptx import Presentation as P

        for mod in [wc, ec, pc]:
            mod.OUTPUT_DIR = _tf.mkdtemp()

        # 11.1 — Word: commercial_lease cambia título por idioma
        fp_es = wc.generate_word({"template_id": "commercial_lease",
            "titulo": "WLeaseES_11", "instrucciones": "",
            "datos": {"idioma": "es", "monto": "8500", "duracion": "24"}})
        fp_en = wc.generate_word({"template_id": "commercial_lease",
            "titulo": "WLeaseEN_11", "instrucciones": "",
            "datos": {"idioma": "en", "monto": "8500", "duracion": "24"}})
        t_es = Document(fp_es).paragraphs[0].text
        t_en = Document(fp_en).paragraphs[0].text
        assert "CONTRATO" in t_es, f"Word ES no en español: '{t_es}'"
        assert "COMMERCIAL" in t_en and "AGREEMENT" in t_en, f"Word EN no en inglés: '{t_en}'"
        print(f"  Word commercial_lease: ES='{t_es[:30]}', EN='{t_en[:30]}'")

        # 11.2 — Excel: sheet names cambian para 5 templates clave
        excel_checks = [
            ("inventory", "Inventario", "Inventory"),
            ("payroll", "Nómina", "Payroll"),
            ("cashflow", "Flujo", "Cash Flow"),
            ("balance_sheet", "Balance General", "Balance Sheet"),
            ("purchase_order", "Orden", "Purchase Order"),
        ]
        for tid, exp_es, exp_en in excel_checks:
            fp_es,_ = ec.generate_excel({"template_id": tid, "titulo": f"X_{tid}_es_11",
                "instrucciones": "", "datos": {"idioma": "es"}})
            fp_en,_ = ec.generate_excel({"template_id": tid, "titulo": f"X_{tid}_en_11",
                "instrucciones": "", "datos": {"idioma": "en"}})
            s_es = load_workbook(fp_es).sheetnames[0]
            s_en = load_workbook(fp_en).sheetnames[0]
            assert exp_es in s_es, f"Excel {tid} ES: esperado '{exp_es}' en '{s_es}'"
            assert exp_en in s_en, f"Excel {tid} EN: esperado '{exp_en}' en '{s_en}'"
        print(f"  Excel: {len(excel_checks)}/{len(excel_checks)} templates traducen sheet names")

        # 11.3 — PowerPoint: al menos 3 planners producen runs distintos en EN vs ES
        for tid in ["pitch", "sales", "report"]:
            fp_es = pc.generate_powerpoint({"template_id": tid, "titulo": f"P_{tid}_es_11",
                "instrucciones": "", "datos": {"idioma": "es"}})
            fp_en = pc.generate_powerpoint({"template_id": tid, "titulo": f"P_{tid}_en_11",
                "instrucciones": "", "datos": {"idioma": "en"}})
            txt_es = [r.text for s in P(fp_es).slides for sh in s.shapes
                      if hasattr(sh, 'text_frame') and sh.has_text_frame
                      for p in sh.text_frame.paragraphs for r in p.runs if r.text]
            txt_en = [r.text for s in P(fp_en).slides for sh in s.shapes
                      if hasattr(sh, 'text_frame') and sh.has_text_frame
                      for p in sh.text_frame.paragraphs for r in p.runs if r.text]
            diffs = sum(1 for a, b in zip(txt_es, txt_en) if a != b and a and b)
            assert diffs >= 3, f"PPT {tid}: solo {diffs} runs traducidos (esperaba ≥3)"
        print(f"  PowerPoint: 3/3 planners producen runs traducidos en EN")
    except ImportError as e:
        print(f"  SKIP ({e})")

    # ── TEST 12 — v14.5: Production-ready (fechas US, moneda, Stripe fallback) ──
    print("\n[12] v14.5: production-ready fixes")
    try:
        import tempfile as _tf
        import os as _os
        from datetime import datetime as _dt
        from office import word_controller as wc
        from office import excel_controller as ec
        from docx import Document
        from openpyxl import load_workbook
        for mod in [wc, ec]:
            mod.OUTPUT_DIR = _tf.mkdtemp()

        # 12.1 — _fmt_short_date respeta idioma
        d = _dt(2026, 5, 21)
        assert wc._fmt_short_date({"idioma":"es"}, d) == "21/05/2026"
        assert wc._fmt_short_date({"idioma":"en"}, d) == "05/21/2026"
        assert wc._fmt_short_date(None, d) == "21/05/2026"  # default ES
        print(f"  _fmt_short_date: ES=21/05/2026, EN=05/21/2026 ✓")

        # 12.2 — Word footer en EN dice "Generated with WEP AI" con formato US
        fp = wc.generate_word({"template_id":"nda","titulo":"NDA_v145_test",
            "instrucciones":"","datos":{"idioma":"en"}})
        all_text = "\n".join(p.text for p in Document(fp).paragraphs)
        assert "Generated with WEP AI" in all_text, "footer EN no traducido"
        assert "Documento generado con" not in all_text, "footer ES sigue presente"
        print(f"  Word footer bilingüe ✓")

        # 12.3 — Excel post-procesador convierte MXN → USD cuando lang=en
        fp,_ = ec.generate_excel({"template_id":"travel_expenses",
            "titulo":"Travel_v145_test","instrucciones":"",
            "datos":{"idioma":"en"}})
        wb = load_workbook(fp)
        all_cells = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for v in row:
                    if isinstance(v, str): all_cells.append(v)
        mxn_count = sum(1 for v in all_cells if 'MXN' in v)
        assert mxn_count == 0, f"MXN debería estar traducido pero hay {mxn_count} celdas"
        print(f"  Excel: MXN convertido a USD ✓")

        # 12.4 — Stripe client: fallback a demo sin claves, detecta config con claves
        from storage import stripe_client
        import importlib
        # Sin claves
        _os.environ.pop("STRIPE_SECRET_KEY", None)
        importlib.reload(stripe_client)
        assert stripe_client.is_configured() == False
        assert stripe_client.create_checkout_session("pro","x@y.com",1) is None
        # Con clave dummy
        _os.environ["STRIPE_SECRET_KEY"] = "sk_test_dummy_v145"
        importlib.reload(stripe_client)
        assert stripe_client.is_configured() == True
        # Limpieza
        _os.environ.pop("STRIPE_SECRET_KEY", None)
        print(f"  Stripe client: detección + fallback demo ✓")
    except ImportError as e:
        print(f"  SKIP ({e})")

    # ── TEST 13 — v14.6: ui_lang fallback + email + trial + billing helpers ──
    print("\n[13] v14.6: ui_lang fallback, email, trial, billing")
    try:
        import os as _os
        from unittest.mock import patch as _patch, MagicMock as _MM
        _os.environ.setdefault("ANTHROPIC_API_KEY", "stub")
        from agent import brain
        from storage import email, stripe_client, db
        import importlib

        def _mock_response(json_payload):
            mr = _MM()
            mr.content = [_MM(text=f"OK\n\n##GENERAR##\n{json_payload}")]
            return mr

        # 13.1 — ui_lang fallback se inyecta sólo si LLM no emite idioma
        with _patch.object(brain, 'client') as mc:
            mc.messages.create.return_value = _mock_response(
                '{"tipo":"word","template_id":"nda","titulo":"X","datos":{}}')
            r = brain.chat([], "Generate", ui_lang="en")
            assert r["gen_data"]["datos"].get("idioma") == "en", \
                "fallback ui_lang=en falló"

            mc.messages.create.return_value = _mock_response(
                '{"tipo":"word","template_id":"nda","titulo":"X","datos":{"idioma":"english"}}')
            r = brain.chat([], "Generate", ui_lang="es")
            assert r["gen_data"]["datos"].get("idioma") == "english", \
                "LLM debe ganar sobre ui_lang"
            print(f"  ui_lang fallback: LLM gana sobre UI cuando emite ✓")

        # 13.2 — email module: is_configured() false sin provider
        _os.environ.pop("EMAIL_PROVIDER", None)
        importlib.reload(email)
        assert email.is_configured() == False
        # send_* devuelve False en dry-run
        assert email.send_welcome("test@example.com", "Test", "en") == False
        assert email.send_payment_confirmed("test@example.com", "pro", "$4.99", "", "es") == False
        print(f"  email module: dry-run sin EMAIL_PROVIDER ✓")

        # 13.3 — stripe_client: trial_days parameter
        _os.environ.pop("STRIPE_SECRET_KEY", None)
        importlib.reload(stripe_client)
        # Sin claves devuelve None igual aunque pase trial_days
        assert stripe_client.create_checkout_session("pro","x@y.com",1,trial_days=14) is None
        print(f"  stripe_client.trial_days: signature acepta ✓")

        # 13.4 — db.get_user_field: whitelist + return None para campo no permitido
        # Crear usuario de prueba
        db.init_db()
        uid_test, err_t = db.register_user("Test","User",
                                            f"v146test_{_os.urandom(4).hex()}@x.com",
                                            "password123","free")
        assert err_t is None, f"register fallo: {err_t}"
        # Campo permitido
        plan = db.get_user_field(uid_test, "plan")
        assert plan == "free", f"plan esperado 'free', got {plan}"
        # Campo NO en whitelist devuelve None
        bad = db.get_user_field(uid_test, "password_hash")
        assert bad is None, f"campo no whitelisted debería ser None, got {bad}"
        # Campo en whitelist pero columna inexistente devuelve None (no excepción)
        stripe_id = db.get_user_field(uid_test, "stripe_customer_id")
        # Puede ser None (columna existe pero vacía) o None (columna no existe)
        # Lo importante es que no lance excepción
        print(f"  db.get_user_field: whitelist + columna inexistente OK ✓")
    except ImportError as e:
        print(f"  SKIP ({e})")
    except Exception as e:
        print(f"  FAIL: {e}")
        raise

    # ── TEST 14 — v14.7: integración email + status banner + cancel direct ──
    print("\n[14] v14.7: email integration, status banner, cancel direct")
    try:
        import os as _os
        from unittest.mock import patch as _patch, MagicMock as _MM
        from storage import email, stripe_client, db
        import importlib

        # 14.1 — Stripe metadata incluye 'lang' y 'email'
        _os.environ["STRIPE_SECRET_KEY"] = "sk_test_v147"
        _os.environ["STRIPE_PRICE_ID_PRO"] = "price_test_v147"
        importlib.reload(stripe_client)

        with _patch('stripe.checkout.Session.create') as mock_create:
            mock_create.return_value = _MM(url="https://checkout.stripe.com/x",
                                            id="cs_x")
            r = stripe_client.create_checkout_session(
                "pro", "test@example.com", 42,
                trial_days=0, user_lang="en")
            # Verificar el metadata pasado a Stripe
            call_kwargs = mock_create.call_args.kwargs
            assert call_kwargs["metadata"]["lang"] == "en"
            assert call_kwargs["metadata"]["email"] == "test@example.com"
            assert call_kwargs["metadata"]["user_id"] == "42"
            assert call_kwargs["subscription_data"]["metadata"]["lang"] == "en"
            print(f"  Stripe metadata incluye lang+email para webhook ✓")
        _os.environ.pop("STRIPE_SECRET_KEY", None)
        _os.environ.pop("STRIPE_PRICE_ID_PRO", None)

        # 14.2 — email.send_welcome devuelve False en dry-run (sin provider)
        # pero NO lanza excepción
        _os.environ.pop("EMAIL_PROVIDER", None)
        importlib.reload(email)
        try:
            ret = email.send_welcome("test@example.com", "Test User", lang="en")
            assert ret == False, "dry-run debería devolver False"
            print(f"  email.send_welcome dry-run no-rompe ✓")
        except Exception as e:
            print(f"  ✗ send_welcome lanzó excepción: {e}")
            raise

        # 14.3 — cancel_subscription helper (sin claves → False)
        _os.environ.pop("STRIPE_SECRET_KEY", None)
        importlib.reload(stripe_client)
        assert stripe_client.cancel_subscription("sub_test_x") == False
        print(f"  stripe_client.cancel_subscription dry-run ✓")

        # 14.4 — db.get_user_field maneja 'subscription_status' aunque no exista la columna
        db.init_db()
        uid_test, err_t = db.register_user("V147","Test",
            f"v147test_{_os.urandom(4).hex()}@x.com", "password123", "free")
        assert err_t is None
        # subscription_status probablemente no existe en DB nueva — debe devolver None sin error
        status = db.get_user_field(uid_test, "subscription_status")
        # Puede ser None (columna existe pero vacía) o None (columna no existe)
        print(f"  db.get_user_field('subscription_status') = {status!r} (sin excepción) ✓")
    except ImportError as e:
        print(f"  SKIP ({e})")
    except Exception as e:
        print(f"  FAIL: {e}")
        raise

    # ── TEST 15 — v14.8: hardening de seguridad + retry + i18n fallback ──────
    print("\n[15] v14.8: subprocess hardening, model env vars, lang fallback")
    try:
        import importlib
        import inspect

        # 15.1 — open_in_word/excel/powerpoint usan `open -a` no `osascript`
        from office import word_controller, excel_controller, ppt_controller
        for mod, name in [(word_controller, "open_in_word"),
                          (excel_controller, "open_in_excel"),
                          (ppt_controller, "open_in_powerpoint")]:
            src = inspect.getsource(getattr(mod, name))
            assert "osascript" not in src, f"{name} todavía usa osascript (riesgo de inyección)"
            assert "open" in src and "-a" in src, f"{name} debería usar subprocess.run(['open','-a',...])"
            assert "FileNotFoundError" in src, f"{name} debería validar existencia del archivo"
        print(f"  3 open_in_* migrados a `open -a` sin AppleScript ✓")

        # 15.2 — brain.py usa constantes MODEL_* en lugar de strings hardcoded
        try:
            from agent import brain
            assert hasattr(brain, "MODEL_CHAT"), "brain.MODEL_CHAT no definido"
            assert hasattr(brain, "MODEL_VISION"), "brain.MODEL_VISION no definido"
            assert hasattr(brain, "MODEL_TITLE"), "brain.MODEL_TITLE no definido"
            # Override por env var debe funcionar
            _os = __import__("os")
            _os.environ["WEPAI_MODEL_CHAT"] = "claude-test-override"
            importlib.reload(brain)
            assert brain.MODEL_CHAT == "claude-test-override", "env var override no funciona"
            _os.environ.pop("WEPAI_MODEL_CHAT", None)
            importlib.reload(brain)
            print(f"  MODEL_CHAT/VISION/TITLE configurables vía env var ✓")
        except ImportError:
            print(f"  SKIP brain.py (anthropic SDK no instalado en sandbox)")

        # 15.3 — _get_user_email_and_lang lee lang real si la columna existe
        # Creamos una DB temporal con columna 'lang' y un user en español
        import tempfile as _tf, sqlite3 as _sq3
        tmp_wh = _tf.NamedTemporaryFile(delete=False, suffix=".db")
        tmp_wh.close()
        _c = _sq3.connect(tmp_wh.name)
        _c.executescript("""
            CREATE TABLE users (
                id INTEGER PRIMARY KEY,
                email TEXT,
                lang TEXT,
                stripe_customer_id TEXT
            );
            INSERT INTO users (email, lang, stripe_customer_id)
            VALUES ('es@test.com', 'es', 'cus_es'),
                   ('en@test.com', 'en', 'cus_en'),
                   ('no_lang@test.com', NULL, 'cus_default');
        """)
        _c.commit(); _c.close()
        try:
            from storage import webhook_server as _wh
            _wh.DATABASE_URL = f"sqlite:///{tmp_wh.name}"
            email_es, lang_es = _wh._get_user_email_and_lang("cus_es")
            email_en, lang_en = _wh._get_user_email_and_lang("cus_en")
            email_def, lang_def = _wh._get_user_email_and_lang("cus_default")
            assert (email_es, lang_es) == ("es@test.com", "es"), f"esperado es@test.com/es, obtuve {email_es!r}/{lang_es!r}"
            assert (email_en, lang_en) == ("en@test.com", "en"), f"esperado en@test.com/en, obtuve {email_en!r}/{lang_en!r}"
            assert lang_def == "en", f"NULL lang debería caer a 'en', obtuve {lang_def!r}"
            print(f"  _get_user_email_and_lang lee lang real de DB (es/en/default) ✓")
        except ImportError:
            print(f"  SKIP webhook_server (flask/stripe no instalados en sandbox)")
        finally:
            os.unlink(tmp_wh.name)
    except Exception as e:
        print(f"  FAIL: {e}")
        raise

    # ── TEST 16 — v14.9: Excel respeta datos del usuario en 5 templates demo ──
    print("\n[16] v14.9: Excel data integration (inventory, quotation, payroll, sales, cashflow)")
    try:
        import tempfile as _tf
        from openpyxl import load_workbook as _lw

        with _tf.TemporaryDirectory() as _tmp:
            os.environ["HOME"] = _tmp
            import importlib as _imp
            from office import excel_controller as _ec
            _imp.reload(_ec)

            # 16.1 — Inventory con productos reales
            gd = {"tipo":"excel","template_id":"inventory","titulo":"Inv Test",
                  "datos":{"items":[
                      {"codigo":"X-1","nombre":"Producto Alpha","stock":100,"precio_costo":50,"precio_venta":100,"categoria":"Beta"},
                      {"codigo":"X-2","nombre":"Producto Gamma","stock":50,"precio_costo":30,"precio_venta":60,"categoria":"Delta"},
                  ]}}
            fp, _ = _ec.generate_excel(gd)
            ws = _lw(fp)["Inventario"]
            cells_text = " ".join(str(ws.cell(r,c).value or "") for r in range(5,7) for c in range(1,11))
            assert "Producto Alpha" in cells_text, f"inventory no integró nombre: {cells_text[:200]}"
            assert "X-1" in cells_text and "X-2" in cells_text, "inventory no integró códigos"
            assert "Beta" in cells_text and "Delta" in cells_text, "inventory no integró categorías"
            print(f"  inventory integra productos del usuario ✓")

            # 16.2 — Quotation con cliente + servicios
            gd = {"tipo":"excel","template_id":"quotation","titulo":"Cot Test",
                  "datos":{"cliente_nombre":"Cliente Demo SA","items":[
                      {"descripcion":"Servicio Alpha","unidad":"Hr","cantidad":10,"precio_venta":150},
                  ]}}
            fp, _ = _ec.generate_excel(gd)
            ws = _lw(fp)["Cotización"]
            cells_text = " ".join(str(ws.cell(r,c).value or "") for r in range(4,13) for c in range(1,8))
            assert "Cliente Demo SA" in cells_text, "quotation no integró cliente"
            assert "Servicio Alpha" in cells_text, "quotation no integró item"
            print(f"  quotation integra cliente + servicios del usuario ✓")

            # 16.3 — Payroll con empleados reales
            gd = {"tipo":"excel","template_id":"payroll","titulo":"Nom Test",
                  "datos":{"items":[
                      {"empleado":"Ana Test","puesto":"CTO","dias":30,"salario_diario":1500},
                      {"empleado":"Luis Test","puesto":"Dev","dias":30,"salario_diario":1000},
                  ]}}
            fp, _ = _ec.generate_excel(gd)
            ws = _lw(fp)["Nómina"]
            cells_text = " ".join(str(ws.cell(r,c).value or "") for r in range(5,8) for c in range(1,12))
            assert "Ana Test" in cells_text and "Luis Test" in cells_text, "payroll no integró empleados"
            assert "CTO" in cells_text, "payroll no integró puestos"
            print(f"  payroll integra empleados del usuario ✓")

            # 16.4 — Sales con ventas reales
            gd = {"tipo":"excel","template_id":"sales","titulo":"Ventas Test",
                  "datos":{"items":[
                      {"fecha":"2026-06-01","producto":"Producto Real","cliente":"Cliente Real","cantidad":1,"precio":499},
                  ]}}
            fp, _ = _ec.generate_excel(gd)
            ws = _lw(fp)["Ventas"]
            cells_text = " ".join(str(ws.cell(r,c).value or "") for r in range(5,7) for c in range(1,9))
            assert "Producto Real" in cells_text, "sales no integró producto"
            assert "Cliente Real" in cells_text, "sales no integró cliente"
            print(f"  sales integra ventas del usuario ✓")

            # 16.5 — Cashflow con ingresos/egresos por mes
            gd = {"tipo":"excel","template_id":"cashflow","titulo":"Flujo Test",
                  "datos":{
                      "ingresos":{"Ventas Reales":{"ene":50000,"feb":60000}},
                      "egresos":{"Costos Reales":{"ene":20000,"feb":22000}},
                  }}
            fp, _ = _ec.generate_excel(gd)
            ws = _lw(fp)["Flujo de Caja"]
            cells_text = " ".join(str(ws.cell(r,c).value or "") for r in range(5,20) for c in range(1,15))
            assert "Ventas Reales" in cells_text, "cashflow no integró concepto de ingresos"
            assert "Costos Reales" in cells_text, "cashflow no integró concepto de egresos"
            assert "50000" in cells_text or "60000" in cells_text, "cashflow no integró montos"
            print(f"  cashflow integra ingresos/egresos del usuario ✓")

            # 16.6 — Backward compat: sin items, sigue funcionando
            gd_no_items = {"tipo":"excel","template_id":"inventory","titulo":"BC Test",
                           "datos":{"idioma":"es"}}
            fp, _ = _ec.generate_excel(gd_no_items)
            ws = _lw(fp)["Inventario"]
            cells_text = " ".join(str(ws.cell(r,c).value or "") for r in range(5,10) for c in range(2,5))
            assert "Producto de ejemplo" in cells_text, "backward compat roto: sin items debería usar sample"
            print(f"  backward compat: sin items cae a sample data ✓")

            # 16.7 — Helper _norm_item y _num
            assert _ec._norm_item({"code":"X","name":"Y","qty":5}) == {"codigo":"X","nombre":"Y","cantidad":5}
            assert _ec._num("$1,234.56") == 1234.56
            assert _ec._num("50%") == 0.5
            assert _ec._num(None, default=42) == 42
            print(f"  helpers _norm_item + _num funcionan correctamente ✓")
    except Exception as e:
        print(f"  FAIL: {e}")
        raise

    print("\n" + "=" * 50)
    print("All 16 test groups PASSED.")


if __name__ == "__main__":
    main()
