"""
WEP AI — PowerPoint Controller v2
Genera presentaciones profesionales .pptx desde lenguaje natural.

8 temas · 12 layouts · 12 deck types · datos reales · gráficos · speaker notes
"""

import subprocess, os, re, io, threading
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime

OUTPUT_DIR = os.path.expanduser("~/Documents/WEP_AI")
W = Inches(13.33)
H = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════════
# v14.4 — I18N (ES/EN) thread-safe
# ══════════════════════════════════════════════════════════════════════════════
_pp_ctx = threading.local()


def _pp_set_lang(lang: str):
    """Setea idioma (thread-local) para esta llamada a generate_powerpoint."""
    if not lang:
        lang = "es"
    s = str(lang).strip().lower()
    if s in ("es", "español", "espanol", "spanish", "esp"):
        _pp_ctx.lang = "es"
    elif s in ("en", "english", "inglés", "ingles", "eng"):
        _pp_ctx.lang = "en"
    else:
        _pp_ctx.lang = "es"


def _pp_lang() -> str:
    return getattr(_pp_ctx, "lang", "es")


# Patrones ES → EN para post-procesamiento de slides. Mantenido en sincronía
# (subset) con los de excel_controller para que ambos productos hablen igual.
_PP_PATTERNS = [
    # ── Frases completas / títulos de slides ────────────────────────────────
    (r"\bResumen Ejecutivo\b",      "Executive Summary"),
    (r"\bAgenda\b",                 "Agenda"),
    (r"\bIntroducción\b",           "Introduction"),
    (r"\bIntroduccion\b",           "Introduction"),
    (r"\bConclusión\b",             "Conclusion"),
    (r"\bConclusiones\b",           "Conclusions"),
    (r"\bConclusion\b",             "Conclusion"),
    (r"\bRecomendaciones\b",        "Recommendations"),
    (r"\bRecomendación\b",          "Recommendation"),
    (r"\bSiguientes pasos\b",       "Next steps"),
    (r"\bPróximos pasos\b",         "Next steps"),
    (r"\bProximos pasos\b",         "Next steps"),
    (r"\bGracias\b",                "Thank you"),
    (r"\bMuchas Gracias\b",         "Thank you"),
    (r"\bPreguntas\b",              "Questions"),
    (r"\bContacto\b",               "Contact"),
    (r"\bContáctanos\b",            "Contact us"),
    (r"\b¿Preguntas\?",             "Questions?"),
    (r"\bSobre Nosotros\b",         "About Us"),
    (r"\bNuestra Empresa\b",        "Our Company"),
    (r"\bNuestro Equipo\b",         "Our Team"),
    (r"\bNuestra Misión\b",         "Our Mission"),
    (r"\bNuestra Visión\b",         "Our Vision"),
    (r"\bMisión\b",                 "Mission"),
    (r"\bVisión\b",                 "Vision"),
    (r"\bValores\b",                "Values"),
    (r"\bEl Problema\b",            "The Problem"),
    (r"\bLa Solución\b",            "The Solution"),
    (r"\bSolución\b",               "Solution"),
    (r"\bMercado Objetivo\b",       "Target Market"),
    (r"\bTamaño de Mercado\b",      "Market Size"),
    (r"\bModelo de Negocio\b",      "Business Model"),
    (r"\bPropuesta de Valor\b",     "Value Proposition"),
    (r"\bVentaja Competitiva\b",    "Competitive Advantage"),
    (r"\bCompetencia\b",            "Competition"),
    (r"\bCompetidores\b",           "Competitors"),
    (r"\bEstrategia\b",             "Strategy"),
    (r"\bRoadmap\b",                "Roadmap"),
    (r"\bHoja de Ruta\b",           "Roadmap"),
    (r"\bCronograma\b",             "Timeline"),
    (r"\bHitos\b",                  "Milestones"),
    (r"\bEntregables\b",            "Deliverables"),
    (r"\bResultados Clave\b",       "Key Results"),
    (r"\bMétricas Clave\b",         "Key Metrics"),
    (r"\bIndicadores Clave\b",      "Key Indicators"),
    (r"\bKPIs Clave\b",             "Key KPIs"),
    (r"\bAnálisis FODA\b",          "SWOT Analysis"),
    (r"\bFODA\b",                   "SWOT"),
    (r"\bFortalezas\b",             "Strengths"),
    (r"\bOportunidades\b",          "Opportunities"),
    (r"\bDebilidades\b",            "Weaknesses"),
    (r"\bAmenazas\b",               "Threats"),
    (r"\bProyecciones Financieras\b","Financial Projections"),
    (r"\bProyección Financiera\b",  "Financial Projection"),
    (r"\bInversión Solicitada\b",   "Funding Ask"),
    (r"\bUso de Fondos\b",          "Use of Funds"),
    (r"\bEquipo Fundador\b",        "Founding Team"),
    (r"\bEquipo\b",                 "Team"),
    (r"\bTracción\b",               "Traction"),
    (r"\bTraccion\b",               "Traction"),
    (r"\bCrecimiento\b",            "Growth"),
    (r"\bIngresos Recurrentes\b",   "Recurring Revenue"),
    (r"\bCAC\b",                    "CAC"),
    (r"\bLTV\b",                    "LTV"),
    (r"\bChurn\b",                  "Churn"),
    (r"\bResultados\b",             "Results"),
    (r"\bResultado\b",              "Result"),
    (r"\bDesafíos\b",               "Challenges"),
    (r"\bDesafios\b",               "Challenges"),
    (r"\bLogros\b",                 "Achievements"),
    (r"\bAprendizajes\b",           "Lessons Learned"),
    (r"\bResumen\b",                "Summary"),
    (r"\bDescripción\b",            "Description"),
    (r"\bDescripcion\b",            "Description"),
    (r"\bObjetivos\b",              "Objectives"),
    (r"\bObjetivo\b",               "Objective"),
    (r"\bMetas\b",                  "Goals"),
    (r"\bMeta\b",                   "Goal"),
    (r"\bAlcance\b",                "Scope"),
    (r"\bFase\b",                   "Phase"),
    (r"\bEtapa\b",                  "Stage"),
    (r"\bSemana\b",                 "Week"),
    (r"\bMes\b",                    "Month"),
    (r"\bAño\b",                    "Year"),
    (r"\bTrimestre\b",              "Quarter"),
    (r"\bPresentado por\b",         "Presented by"),
    (r"\bPreparado por\b",          "Prepared by"),
    (r"\bConfidencial\b",           "Confidential"),
    (r"\bAcerca de\b",              "About"),
    (r"\bCaso de Éxito\b",          "Success Story"),
    (r"\bCaso de Exito\b",          "Success Story"),
    (r"\bCaso de Estudio\b",        "Case Study"),
    (r"\bAntes\b",                  "Before"),
    (r"\bDespués\b",                "After"),
    (r"\bDespues\b",                "After"),
    (r"\bImpacto\b",                "Impact"),
    (r"\bBeneficios\b",             "Benefits"),
    (r"\bBeneficio\b",              "Benefit"),
    (r"\bCaracterísticas\b",        "Features"),
    (r"\bCaracteristicas\b",        "Features"),
    (r"\bFuncionalidades\b",        "Features"),
    (r"\bPricing\b",                "Pricing"),
    (r"\bPrecios\b",                "Pricing"),
    (r"\bPlan\b",                   "Plan"),
    (r"\bPlanes\b",                 "Plans"),
    (r"\bGratis\b",                 "Free"),
    (r"\bMensual\b",                "Monthly"),
    (r"\bAnual\b",                  "Annual"),
    (r"\bPropuesta\b",              "Proposal"),
    (r"\bDemo\b",                   "Demo"),
    (r"\bConócenos\b",              "Meet us"),
    (r"\bQuiénes Somos\b",          "Who We Are"),
    (r"\bQuienes Somos\b",          "Who We Are"),
    (r"\bQué Hacemos\b",            "What We Do"),
    (r"\bQue Hacemos\b",            "What We Do"),
    (r"\bPor Qué Nosotros\b",       "Why Us"),
    (r"\bPor Que Nosotros\b",       "Why Us"),
    (r"\bMercado\b",                "Market"),
    (r"\bIndustria\b",              "Industry"),
    (r"\bSector\b",                 "Sector"),
    (r"\bClientes\b",               "Clients"),
    (r"\bUsuarios\b",               "Users"),
    (r"\bPlataforma\b",             "Platform"),
    (r"\bSeguridad\b",              "Security"),
    (r"\bSoporte\b",                "Support"),
    (r"\bIntegraciones\b",          "Integrations"),
    (r"\bCasos de Uso\b",           "Use Cases"),
    (r"\bTestimonios\b",            "Testimonials"),
    (r"\bReferencias\b",            "References"),
    (r"\bCertificaciones\b",        "Certifications"),
    (r"\bPartners\b",               "Partners"),
    (r"\bSocios\b",                 "Partners"),
    (r"\bAlianzas\b",               "Partnerships"),
    (r"\bContexto\b",               "Context"),
    (r"\bSituación Actual\b",       "Current Situation"),
    (r"\bSituacion Actual\b",       "Current Situation"),
    (r"\bOportunidad\b",            "Opportunity"),
    (r"\bRiesgos\b",                "Risks"),
    (r"\bRiesgo\b",                 "Risk"),
    (r"\bMitigación\b",             "Mitigation"),
    (r"\bMitigacion\b",             "Mitigation"),
    (r"\bSupuestos\b",              "Assumptions"),
    (r"\bDependencias\b",           "Dependencies"),
    (r"\bAcciones\b",               "Actions"),
    (r"\bAcción\b",                 "Action"),
    (r"\bAccion\b",                 "Action"),
    (r"\bResponsables\b",           "Owners"),
    (r"\bResponsable\b",            "Owner"),
    (r"\bFechas Clave\b",           "Key Dates"),
    (r"\bFecha\b",                  "Date"),
    (r"\bTotal\b",                  "Total"),
    (r"\bSubtotal\b",               "Subtotal"),
    (r"\bGracias por su tiempo\b",  "Thank you for your time"),
    (r"\bGracias por su atención\b","Thank you for your attention"),
    (r"\bGracias por escucharnos\b","Thank you for listening"),
    # Meses (forma corta y larga)
    (r"\bEnero\b",     "January"),  (r"\bFebrero\b",    "February"),
    (r"\bMarzo\b",     "March"),    (r"\bAbril\b",      "April"),
    (r"\bMayo\b",      "May"),      (r"\bJunio\b",      "June"),
    (r"\bJulio\b",     "July"),     (r"\bAgosto\b",     "August"),
    (r"\bSeptiembre\b","September"),(r"\bOctubre\b",    "October"),
    (r"\bNoviembre\b", "November"), (r"\bDiciembre\b",  "December"),
    (r"\bEne\b","Jan"),(r"\bFeb\b","Feb"),(r"\bMar\b","Mar"),
    (r"\bAbr\b","Apr"),(r"\bMay\b","May"),(r"\bJun\b","Jun"),
    (r"\bJul\b","Jul"),(r"\bAgo\b","Aug"),(r"\bSep\b","Sep"),
    (r"\bOct\b","Oct"),(r"\bNov\b","Nov"),(r"\bDic\b","Dec"),
    (r"\bde de\b",                  "of"),   # fix doble preposición tras fecha
    (r"(\d+) de (\w+) de (\d{4})",  r"\2 \1, \3"),  # "21 de mayo de 2026" → "mayo 21, 2026"
    (r"\bNuestra Solución\b",       "Our Solution"),
    (r"\bNuestra Solucion\b",       "Our Solution"),
    (r"\bNuestra Solution\b",       "Our Solution"),
    (r"\bNuestro Producto\b",       "Our Product"),
    (r"\bNuestro Product\b",        "Our Product"),
    (r"\bNuestros Productos\b",     "Our Products"),
    (r"\bNuestros Products\b",      "Our Products"),
    (r"\bNuestros Clientes\b",      "Our Clients"),
    (r"\bNuestros Clients\b",       "Our Clients"),
    (r"\bNuestra Propuesta\b",      "Our Proposal"),
    (r"\bNuestra Proposal\b",       "Our Proposal"),
    (r"\bNuestro Enfoque\b",        "Our Approach"),
    (r"\bNuestro Diferencial\b",    "Our Differentiator"),
    (r"\bNuestro Diferentiator\b",  "Our Differentiator"),
    # ── Placeholders entre corchetes ──────────────────────────────────────
    (r"\[Empresa\]",                "[Company]"),
    (r"\[Cliente\]",                "[Client]"),
    (r"\[Producto\]",               "[Product]"),
    (r"\[Sector\]",                 "[Sector]"),
    (r"\[Mercado\]",                "[Market]"),
    (r"\[Mercado objetivo\]",       "[Target Market]"),
    (r"\[Período\]",                "[Period]"),
    (r"\[Periodo\]",                "[Period]"),
    (r"\[Fundador\]",               "[Founder]"),
    (r"\[Año\]",                    "[Year]"),
    (r"\[Instructor\]",             "[Instructor]"),
    (r"\[Monto\]",                  "[Amount]"),
    (r"\[Duración\]",               "[Duration]"),
    (r"\[Duracion\]",               "[Duration]"),
    (r"\[Nombre\]",                 "[Name]"),
    (r"\[Cargo\]",                  "[Title]"),
    (r"\[Departamento\]",           "[Department]"),
    (r"\[Equipo\]",                 "[Team]"),
    (r"\[Fecha\]",                  "[Date]"),
    (r"\[Ciudad\]",                 "[City]"),
    (r"\[País\]",                   "[Country]"),
    (r"\[Pais\]",                   "[Country]"),
    (r"\[Tutor\]",                  "[Advisor]"),
    (r"\[PM\]",                     "[PM]"),
    # Fallback de la closing slide (CTA por defecto del controller)
    (r"Quedamos a su disposición para cualquier consulta\.?",
        "We remain available for any inquiries."),
]
_PP_COMPILED = [(re.compile(p), repl) for p, repl in _PP_PATTERNS]


def _pp_translate(s: str) -> str:
    """Aplica patrones a una cadena. Devuelve traducida o la original."""
    if not isinstance(s, str) or not s.strip():
        return s
    out = s
    for rx, repl in _PP_COMPILED:
        out = rx.sub(repl, out)
    return out


def _translate_presentation(prs):
    """Post-procesa la presentación traduciendo todos los runs de texto.
    Sólo se ejecuta si _pp_lang() == 'en'. Recorre slides → shapes →
    text_frames → paragraphs → runs y también notes_slides."""
    if _pp_lang() != "en":
        return
    for slide in prs.slides:
        # Shapes con texto
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text and isinstance(run.text, str):
                                new = _pp_translate(run.text)
                                if new != run.text:
                                    run.text = new
            except Exception:
                pass
            # Tablas dentro de shapes
            try:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text and isinstance(run.text, str):
                                        new = _pp_translate(run.text)
                                        if new != run.text:
                                            run.text = new
            except Exception:
                pass
            # Charts dentro de shapes — título y axes
            try:
                if shape.has_chart:
                    chart = shape.chart
                    if chart.has_title:
                        title_frame = chart.chart_title.text_frame
                        for paragraph in title_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text:
                                    new = _pp_translate(run.text)
                                    if new != run.text:
                                        run.text = new
            except Exception:
                pass
        # Speaker notes
        try:
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                for paragraph in notes_tf.paragraphs:
                    for run in paragraph.runs:
                        if run.text and isinstance(run.text, str):
                            new = _pp_translate(run.text)
                            if new != run.text:
                                run.text = new
        except Exception:
            pass


# ══════════════════════════════════════════════════════════════════════════════
# SISTEMA DE COLORES — 8 TEMAS
# ══════════════════════════════════════════════════════════════════════════════

def _rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

THEMES = {
    "profesional": {"dark_bg":_rgb("0D1B2A"),"accent":_rgb("185FA5"),"dark_text":_rgb("FFFFFF"),"light_bg":_rgb("F0F4F8"),"light_text":_rgb("0D1B2A"),"muted":_rgb("93C5FD"),"highlight":_rgb("3B82F6"),"font_h":"Calibri","font_b":"Calibri Light"},
    "ejecutivo":   {"dark_bg":_rgb("1A1A2E"),"accent":_rgb("F5A623"),"dark_text":_rgb("FFFFFF"),"light_bg":_rgb("FAFAFA"),"light_text":_rgb("1A1A2E"),"muted":_rgb("FCD34D"),"highlight":_rgb("F59E0B"),"font_h":"Georgia","font_b":"Calibri"},
    "formal":      {"dark_bg":_rgb("0F2027"),"accent":_rgb("00B09B"),"dark_text":_rgb("FFFFFF"),"light_bg":_rgb("F0FDF9"),"light_text":_rgb("064E3B"),"muted":_rgb("6EE7B7"),"highlight":_rgb("10B981"),"font_h":"Calibri","font_b":"Calibri Light"},
    "creativo":    {"dark_bg":_rgb("2E1065"),"accent":_rgb("EC4899"),"dark_text":_rgb("FFFFFF"),"light_bg":_rgb("FDF2F8"),"light_text":_rgb("4A0472"),"muted":_rgb("F9A8D4"),"highlight":_rgb("DB2777"),"font_h":"Trebuchet MS","font_b":"Arial"},
    "startup":     {"dark_bg":_rgb("0F172A"),"accent":_rgb("6366F1"),"dark_text":_rgb("FFFFFF"),"light_bg":_rgb("F8FAFC"),"light_text":_rgb("0F172A"),"muted":_rgb("A5B4FC"),"highlight":_rgb("818CF8"),"font_h":"Arial","font_b":"Arial"},
    "consultora":  {"dark_bg":_rgb("1E3A5F"),"accent":_rgb("C0392B"),"dark_text":_rgb("FFFFFF"),"light_bg":_rgb("FEFEFE"),"light_text":_rgb("1E3A5F"),"muted":_rgb("F1948A"),"highlight":_rgb("E74C3C"),"font_h":"Cambria","font_b":"Calibri"},
    "minimalista": {"dark_bg":_rgb("212121"),"accent":_rgb("424242"),"dark_text":_rgb("FFFFFF"),"light_bg":_rgb("FFFFFF"),"light_text":_rgb("212121"),"muted":_rgb("9E9E9E"),"highlight":_rgb("616161"),"font_h":"Arial Black","font_b":"Arial"},
    "impacto":     {"dark_bg":_rgb("B91C1C"),"accent":_rgb("FCD34D"),"dark_text":_rgb("FFFFFF"),"light_bg":_rgb("FFFBEB"),"light_text":_rgb("7F1D1D"),"muted":_rgb("FCA5A5"),"highlight":_rgb("EF4444"),"font_h":"Impact","font_b":"Arial"},
}

THEME_ALIASES = {
    "azul":"profesional","navy":"profesional","corporativo":"profesional",
    "dorado":"ejecutivo","gold":"ejecutivo","premium":"ejecutivo",
    "verde":"formal","green":"formal","esmeralda":"formal",
    "rosa":"creativo","fucsia":"creativo","morado":"creativo","purple":"creativo",
    "indigo":"startup","tech":"startup","moderno":"startup",
    "rojo":"consultora","red":"consultora","bold":"consultora",
    "blanco":"minimalista","clean":"minimalista","simple":"minimalista","white":"minimalista",
    "naranja":"impacto","energia":"impacto","dinamico":"impacto","orange":"impacto",
}

def _resolve_theme(estilo):
    k = estilo.lower().strip()
    return THEMES.get(THEME_ALIASES.get(k, k), THEMES["profesional"])


# ══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def _ensure_dir():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

def _safe_filename(t):
    return re.sub(r'[^\w\s-]','',t).strip().replace(' ','_')[:60]

def _d(datos, key, default=""):
    return str(datos.get(key) or default).strip()

def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for sh in list(slide.shapes):
        if sh.shape_type == 14:
            sh._element.getparent().remove(sh._element)
    return slide

def _rect(slide, x, y, w, h, fill, line=None, lw=0):
    from pptx.util import Pt as P
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = P(lw)
    else: sh.line.fill.background()
    return sh

def _txt(slide, text, x, y, w, h, size=18, bold=False, italic=False,
         color=None, align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    if color is None: color = _rgb("FFFFFF")
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb

def _notes(slide, text):
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception as e:
        print(f"[WEP AI ppt] _notes failed: {e}")

def _img(slide, buf, x, y, w):
    try:
        return slide.shapes.add_picture(buf, x, y, w)
    except Exception as e:
        print(f"[WEP AI ppt] _img failed: {e}")
        return None

def _insert_logo(slide, logo_path_or_buf, position="bottom_left",
                 max_w_in=1.8, max_h_in=0.7):
    """Inserta el logo del usuario en un slide con tamaño controlado.
    position: 'bottom_left' | 'bottom_right' | 'top_right' | 'center_bottom'
    """
    from pptx.util import Inches as I
    try:
        # Detect format and load
        if isinstance(logo_path_or_buf, str):
            if not os.path.exists(logo_path_or_buf):
                return None
            logo_path_or_buf = logo_path_or_buf

        # Position presets (x, y, max_w, max_h)
        margin = I(0.25)
        presets = {
            "bottom_left":    (margin, H - I(max_h_in) - margin),
            "bottom_right":   (W - I(max_w_in) - margin, H - I(max_h_in) - margin),
            "top_right":      (W - I(max_w_in) - margin, margin),
            "center_bottom":  (W/2 - I(max_w_in)/2, H - I(max_h_in) - margin),
        }
        x, y = presets.get(position, presets["bottom_left"])

        pic = slide.shapes.add_picture(
            logo_path_or_buf, x, y,
            width=I(max_w_in), height=I(max_h_in)
        )
        return pic
    except Exception:
        return None   # Falla silenciosa — nunca rompe el deck


def _chart_bars(labels, values, title="", color="#185FA5", bg="F0F4F8", w=5.5, h=3.2):
    import matplotlib; matplotlib.use('Agg')
    import matplotlib.pyplot as plt, numpy as np
    bh = f"#{bg}"; fig, ax = plt.subplots(figsize=(w, h), dpi=120)
    fig.patch.set_facecolor(bh); ax.set_facecolor(bh)
    xs = np.arange(len(labels))
    bars = ax.bar(xs, values, color=color, width=0.55, edgecolor='white', linewidth=0.8, zorder=3)
    for bar, v in zip(bars, values):
        lbl = f"{v:.0f}%" if max(values)<=100 else f"{v:,.0f}"
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+max(values)*0.018,
                lbl, ha='center', va='bottom', fontsize=9, fontweight='bold', color='#111827')
    ax.set_xticks(xs); ax.set_xticklabels(labels, fontsize=9.5, color='#374151')
    ax.set_title(title, fontsize=11, fontweight='bold', color='#111827', pad=8)
    ax.tick_params(axis='y', labelsize=8, labelcolor='#6B7280')
    for sp in ['top','right']: ax.spines[sp].set_visible(False)
    ax.spines['left'].set_color('#D1D5DB'); ax.spines['bottom'].set_color('#D1D5DB')
    ax.yaxis.grid(True, color='#E5E7EB', linewidth=0.6, zorder=0); ax.set_axisbelow(True)
    plt.tight_layout(pad=0.6)
    buf = io.BytesIO(); fig.savefig(buf, format='png', bbox_inches='tight', dpi=120, facecolor=bh, edgecolor='none'); plt.close(fig); buf.seek(0)
    return buf

def _chart_lines(labels, series, title="", bg="F0F4F8", w=5.5, h=3.2):
    import matplotlib; matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    COLS = ['#185FA5','#0F6E56','#B91C1C','#854F0B','#534AB7']
    bh = f"#{bg}"; fig, ax = plt.subplots(figsize=(w, h), dpi=120)
    fig.patch.set_facecolor(bh); ax.set_facecolor(bh)
    for i, (nm, vals) in enumerate(series.items()):
        ax.plot(labels, vals, color=COLS[i%len(COLS)], linewidth=2.2, marker='o', markersize=5, label=nm, zorder=3)
    ax.set_title(title, fontsize=11, fontweight='bold', color='#111827', pad=8)
    ax.tick_params(labelsize=8, labelcolor='#6B7280')
    for sp in ['top','right']: ax.spines[sp].set_visible(False)
    ax.spines['left'].set_color('#D1D5DB'); ax.spines['bottom'].set_color('#D1D5DB')
    ax.yaxis.grid(True, color='#E5E7EB', linewidth=0.6, zorder=0); ax.set_axisbelow(True)
    if len(series)>1: ax.legend(fontsize=8, framealpha=0.9)
    plt.tight_layout(pad=0.6)
    buf = io.BytesIO(); fig.savefig(buf, format='png', bbox_inches='tight', dpi=120, facecolor=bh, edgecolor='none'); plt.close(fig); buf.seek(0)
    return buf


# ══════════════════════════════════════════════════════════════════════════════
# 12 LAYOUTS DE SLIDE
# ══════════════════════════════════════════════════════════════════════════════

def _L_cover(prs, titulo, subtitulo, autor, fecha, T, logo=None):
    """Portada — fondo oscuro, título grande, franja de acento."""
    slide = _new_slide(prs); fh, fb = T["font_h"], T["font_b"]
    _rect(slide, 0, 0, W, H, T["dark_bg"])
    _rect(slide, 0, H-Inches(1.55), W, Inches(1.55), T["accent"])
    _rect(slide, 0, 0, Inches(0.14), H, T["highlight"])
    _txt(slide, titulo, Inches(0.55), Inches(1.5), Inches(12.2), Inches(2.8),
         size=46, bold=True, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fh)
    _rect(slide, Inches(4.2), Inches(4.4), Inches(4.9), Inches(0.04), T["muted"])
    if subtitulo:
        _txt(slide, subtitulo, Inches(0.55), Inches(4.6), Inches(12.2), Inches(0.85),
             size=22, color=T["muted"], align=PP_ALIGN.CENTER, font=fb)
    footer = f"{autor}  ·  {fecha}" if autor else f"WEP AI  ·  {fecha}"
    _txt(slide, footer, Inches(0.4), H-Inches(1.32), Inches(12.5), Inches(0.55),
         size=14, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fb)
    # Logo del usuario en portada
    if logo:
        _insert_logo(slide, logo, position='bottom_right', max_w_in=1.8, max_h_in=0.65)
    _notes(slide, f"Portada: {titulo}\n{subtitulo}")
    return slide


def _L_divider(prs, numero, titulo_sec, descripcion, T):
    """Divisor de sección — número grande decorativo."""
    slide = _new_slide(prs); fh, fb = T["font_h"], T["font_b"]
    _rect(slide, 0, 0, W, H, T["dark_bg"])
    _rect(slide, W-Inches(0.18), 0, Inches(0.18), H, T["accent"])
    if numero:
        _txt(slide, str(numero), Inches(0.4), Inches(0.3), Inches(3), Inches(4.5),
             size=120, bold=True, color=T["highlight"], align=PP_ALIGN.LEFT, font=fh)
    _txt(slide, titulo_sec, Inches(0.45), Inches(3.6), Inches(10), Inches(1.9),
         size=40, bold=True, color=T["dark_text"], align=PP_ALIGN.LEFT, font=fh)
    if descripcion:
        _txt(slide, descripcion, Inches(0.45), Inches(5.55), Inches(10), Inches(1.1),
             size=20, color=T["muted"], align=PP_ALIGN.LEFT, font=fb)
    _notes(slide, f"Sección {numero}: {titulo_sec}")
    return slide


def _L_bullets(prs, titulo, bullets, T, num="", icon=""):
    """Bullets — header oscuro, contenido claro con jerarquía."""
    slide = _new_slide(prs); fh, fb = T["font_h"], T["font_b"]
    _rect(slide, 0, 0, W, H, T["light_bg"])
    _rect(slide, 0, 0, W, Inches(1.12), T["dark_bg"])
    _rect(slide, 0, Inches(1.12), W, Inches(0.048), T["accent"])
    title_str = (f"{icon}  {titulo}" if icon else titulo)
    _txt(slide, title_str, Inches(0.4), Inches(0.11), Inches(11.5), Inches(0.92),
         size=26, bold=True, color=T["dark_text"], align=PP_ALIGN.LEFT, font=fh)
    if num:
        _txt(slide, num, Inches(12.1), Inches(0.11), Inches(1.0), Inches(0.92),
             size=13, color=T["muted"], align=PP_ALIGN.RIGHT, font=fb)
    y = Inches(1.32)
    for b in bullets:
        if not b.strip(): y += Inches(0.12); continue
        if b.startswith("##"):
            _txt(slide, b.replace("##","").strip(), Inches(0.38), y, Inches(12.5), Inches(0.5),
                 size=18.5, bold=True, color=T["accent"], font=fh, wrap=False)
            y += Inches(0.54)
        elif b.startswith("•") or b.startswith("-"):
            _txt(slide, "  " + b.lstrip("•- ").strip(), Inches(0.65), y, Inches(12.1), Inches(0.46),
                 size=15.5, color=T["light_text"], font=fb)
            y += Inches(0.48)
        else:
            _txt(slide, b, Inches(0.38), y, Inches(12.5), Inches(0.52),
                 size=17, color=T["light_text"], font=fb)
            y += Inches(0.55)
        if y > Inches(6.85): break
    _notes(slide, f"{titulo}\n" + "\n".join(bullets))
    return slide


def _L_two_col(prs, titulo, left_title, left_items, right_title, right_items, T, num=""):
    """Dos columnas simétricas."""
    slide = _new_slide(prs); fh, fb = T["font_h"], T["font_b"]
    _rect(slide, 0, 0, W, H, T["light_bg"])
    _rect(slide, 0, 0, W, Inches(1.12), T["dark_bg"])
    _rect(slide, 0, Inches(1.12), W, Inches(0.048), T["accent"])
    _txt(slide, titulo, Inches(0.4), Inches(0.11), Inches(11.5), Inches(0.92),
         size=26, bold=True, color=T["dark_text"], align=PP_ALIGN.LEFT, font=fh)
    if num:
        _txt(slide, num, Inches(12.1), Inches(0.11), Inches(1.0), Inches(0.92),
             size=13, color=T["muted"], align=PP_ALIGN.RIGHT, font=fb)
    # Divisor
    _rect(slide, Inches(6.62), Inches(1.28), Inches(0.048), Inches(5.92), T["accent"])
    for col_x, ct, ci in [(Inches(0.35), left_title, left_items),(Inches(6.78), right_title, right_items)]:
        cw = Inches(5.92)
        _rect(slide, col_x, Inches(1.28), cw, Inches(0.5), T["accent"])
        _txt(slide, ct, col_x+Inches(0.1), Inches(1.3), cw-Inches(0.2), Inches(0.46),
             size=16, bold=True, color=T["dark_text"], align=PP_ALIGN.LEFT, font=fh, wrap=False)
        y = Inches(1.9)
        for it in ci:
            if not it.strip(): continue
            is_h = it.startswith("##"); txt = it.replace("##","").strip() if is_h else it.lstrip("•- ").strip()
            _txt(slide, ("" if is_h else "›  ")+txt, col_x+Inches(0.1), y, cw-Inches(0.2), Inches(0.5),
                 size=15 if is_h else 14, bold=is_h,
                 color=T["accent"] if is_h else T["light_text"],
                 font=fh if is_h else fb)
            y += Inches(0.48 if is_h else 0.46)
            if y > Inches(7.0): break
    _notes(slide, f"{titulo}\n{left_title}: {left_items}\n{right_title}: {right_items}")
    return slide


def _L_stats(prs, titulo, stats, T, num=""):
    """Stats callouts — valores grandes en cards. stats=[(val,lbl,desc),...]"""
    slide = _new_slide(prs); fh, fb = T["font_h"], T["font_b"]
    _rect(slide, 0, 0, W, H, T["dark_bg"])
    _rect(slide, 0, 0, W, Inches(1.1), T["accent"])
    _txt(slide, titulo, Inches(0.4), Inches(0.1), Inches(11.5), Inches(0.9),
         size=26, bold=True, color=T["dark_text"], align=PP_ALIGN.LEFT, font=fh)
    if num:
        _txt(slide, num, Inches(12.1), Inches(0.1), Inches(1.0), Inches(0.9),
             size=13, color=T["muted"], align=PP_ALIGN.RIGHT, font=fb)
    n = min(len(stats), 4)
    if n == 0: return slide
    cw = Inches(12.4/n); gap = Inches(0.13); ch = Inches(4.9); sx = Inches(0.42); cy = Inches(1.38)
    for i, st in enumerate(stats[:n]):
        v, lbl, desc = (st[0] if len(st)>0 else "—"), (st[1] if len(st)>1 else ""), (st[2] if len(st)>2 else "")
        cx = sx + i*(cw+gap)
        _rect(slide, cx, cy, cw-gap, ch, T["highlight"])
        _txt(slide, str(v), cx+Inches(0.08), cy+Inches(0.28), cw-gap-Inches(0.16), Inches(2.2),
             size=62, bold=True, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fh)
        _txt(slide, lbl, cx+Inches(0.08), cy+Inches(2.55), cw-gap-Inches(0.16), Inches(0.9),
             size=17, bold=True, color=T["muted"], align=PP_ALIGN.CENTER, font=fb)
        if desc:
            _txt(slide, desc, cx+Inches(0.08), cy+Inches(3.5), cw-gap-Inches(0.16), Inches(1.1),
                 size=12.5, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fb)
    _notes(slide, f"{titulo}\n" + " | ".join(f"{s[0]} {s[1]}" for s in stats[:n]))
    return slide


def _L_timeline(prs, titulo, steps, T, num=""):
    """Timeline horizontal — hasta 5 pasos. steps=[(n,titulo,desc),...]"""
    slide = _new_slide(prs); fh, fb = T["font_h"], T["font_b"]
    _rect(slide, 0, 0, W, H, T["light_bg"])
    _rect(slide, 0, 0, W, Inches(1.12), T["dark_bg"])
    _rect(slide, 0, Inches(1.12), W, Inches(0.048), T["accent"])
    _txt(slide, titulo, Inches(0.4), Inches(0.11), Inches(11.5), Inches(0.92),
         size=26, bold=True, color=T["dark_text"], align=PP_ALIGN.LEFT, font=fh)
    if num:
        _txt(slide, num, Inches(12.1), Inches(0.11), Inches(1.0), Inches(0.92),
             size=13, color=T["muted"], align=PP_ALIGN.RIGHT, font=fb)
    n = min(len(steps), 5)
    if n == 0: return slide
    line_y = Inches(3.5)
    _rect(slide, Inches(0.5), line_y, Inches(12.3), Inches(0.055), T["accent"])
    sw = Inches(12.3/n); sx = Inches(0.5)
    for i, st in enumerate(steps[:n]):
        sn, stitle, sdesc = (str(st[0]) if len(st)>0 else str(i+1)), (st[1] if len(st)>1 else f"Paso {i+1}"), (st[2] if len(st)>2 else "")
        cx = sx + i*sw + sw/2; r = Inches(0.42)
        _rect(slide, cx-r, line_y-r, r*2, r*2, T["accent"])
        _txt(slide, sn, cx-r, line_y-r, r*2, r*2, size=20, bold=True, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fh)
        if i%2==0:
            _txt(slide, stitle, cx-sw/2+Inches(0.1), Inches(1.4), sw-Inches(0.2), Inches(1.8), size=15, bold=True, color=T["light_text"], align=PP_ALIGN.CENTER, font=fh)
            if sdesc: _txt(slide, sdesc, cx-sw/2+Inches(0.1), Inches(4.2), sw-Inches(0.2), Inches(1.5), size=12, color=T["light_text"], align=PP_ALIGN.CENTER, font=fb)
        else:
            _txt(slide, stitle, cx-sw/2+Inches(0.1), Inches(4.2), sw-Inches(0.2), Inches(1.8), size=15, bold=True, color=T["light_text"], align=PP_ALIGN.CENTER, font=fh)
            if sdesc: _txt(slide, sdesc, cx-sw/2+Inches(0.1), Inches(1.4), sw-Inches(0.2), Inches(1.5), size=12, color=T["light_text"], align=PP_ALIGN.CENTER, font=fb)
    _notes(slide, f"{titulo}\n" + " → ".join(s[1] for s in steps[:n]))
    return slide


def _L_quote(prs, quote, author, role, T):
    """Cita / testimonio — oscuro, texto centrado grande."""
    slide = _new_slide(prs); fh, fb = T["font_h"], T["font_b"]
    _rect(slide, 0, 0, W, H, T["dark_bg"])
    _rect(slide, 0, 0, Inches(0.38), H, T["accent"])
    _rect(slide, W-Inches(0.38), 0, Inches(0.38), H, T["accent"])
    _txt(slide, "\u201c", Inches(0.75), Inches(0.2), Inches(3), Inches(2.5), size=110, bold=True, color=T["highlight"], align=PP_ALIGN.LEFT, font=fh)
    _txt(slide, quote, Inches(0.75), Inches(1.4), Inches(11.8), Inches(4.0), size=28, italic=True, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fh)
    _rect(slide, Inches(5.0), Inches(5.65), Inches(3.3), Inches(0.038), T["muted"])
    _txt(slide, author, Inches(0.75), Inches(5.82), Inches(11.8), Inches(0.6), size=18, bold=True, color=T["muted"], align=PP_ALIGN.CENTER, font=fh)
    if role: _txt(slide, role, Inches(0.75), Inches(6.45), Inches(11.8), Inches(0.6), size=14, color=T["muted"], align=PP_ALIGN.CENTER, font=fb)
    _notes(slide, f'"{quote}" — {author}, {role}')
    return slide


def _L_agenda(prs, items, titulo_agenda, T):
    """Agenda — columna izquierda oscura, items a la derecha. items=[(n,titulo,tiempo),...]"""
    slide = _new_slide(prs); fh, fb = T["font_h"], T["font_b"]
    _rect(slide, 0, 0, W, H, T["dark_bg"])
    _rect(slide, 0, 0, Inches(4.7), H, T["accent"])
    _txt(slide, titulo_agenda or "AGENDA", Inches(0.3), Inches(2.4), Inches(4.1), Inches(2.2),
         size=38, bold=True, color=T["dark_text"], align=PP_ALIGN.LEFT, font=fh)
    _txt(slide, "WEP AI", Inches(0.3), Inches(6.75), Inches(4.1), Inches(0.5),
         size=13, color=T["dark_text"], align=PP_ALIGN.LEFT, font=fb)
    y = Inches(0.9)
    for it in items[:8]:
        n = str(it[0]) if len(it)>0 else ""
        t = it[1] if len(it)>1 else ""; tm = it[2] if len(it)>2 else ""
        _rect(slide, Inches(5.1), y, Inches(0.55), Inches(0.52), T["accent"])
        _txt(slide, n, Inches(5.1), y, Inches(0.55), Inches(0.52), size=16, bold=True, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fh)
        _txt(slide, t, Inches(5.78), y, Inches(6.0), Inches(0.52), size=18, bold=True, color=T["dark_text"], align=PP_ALIGN.LEFT, font=fh)
        if tm: _txt(slide, tm, Inches(11.85), y, Inches(1.25), Inches(0.52), size=14, color=T["muted"], align=PP_ALIGN.RIGHT, font=fb)
        y += Inches(0.73)
        if y > Inches(6.9): break
    _notes(slide, "Agenda:\n" + "\n".join(f"{it[0]}. {it[1]}" for it in items))
    return slide


def _L_chart(prs, titulo, chart_buf, caption, bullets, T, num=""):
    """Gráfico a la izquierda, bullets a la derecha."""
    slide = _new_slide(prs); fh, fb = T["font_h"], T["font_b"]
    _rect(slide, 0, 0, W, H, T["light_bg"])
    _rect(slide, 0, 0, W, Inches(1.12), T["dark_bg"])
    _rect(slide, 0, Inches(1.12), W, Inches(0.048), T["accent"])
    _txt(slide, titulo, Inches(0.4), Inches(0.11), Inches(11.5), Inches(0.92),
         size=26, bold=True, color=T["dark_text"], align=PP_ALIGN.LEFT, font=fh)
    if num:
        _txt(slide, num, Inches(12.1), Inches(0.11), Inches(1.0), Inches(0.92),
             size=13, color=T["muted"], align=PP_ALIGN.RIGHT, font=fb)
    if chart_buf:
        try:
            _img(slide, chart_buf, Inches(0.3), Inches(1.25), Inches(7.6))
        except Exception as e:
            print(f"[WEP AI ppt] chart insert failed: {e}")
    if caption: _txt(slide, caption, Inches(0.3), Inches(6.55), Inches(7.6), Inches(0.55), size=11, italic=True, color=T["muted"], align=PP_ALIGN.CENTER, font=fb)
    y = Inches(1.4)
    for b in bullets[:7]:
        if not b.strip(): continue
        is_h = b.startswith("##"); txt = b.replace("##","").strip() if is_h else b.lstrip("•- ")
        _txt(slide, ("" if is_h else "›  ")+txt, Inches(8.15), y, Inches(4.9), Inches(0.55),
             size=15 if is_h else 14, bold=is_h, color=T["accent"] if is_h else T["light_text"], font=fh if is_h else fb)
        y += Inches(0.58); 
        if y > Inches(6.9): break
    _notes(slide, f"{titulo}\n{caption}")
    return slide


def _L_team(prs, titulo, members, T, num=""):
    """Team grid — cards con avatar, nombre, rol, descripción. members=[(nombre,rol,desc),...]"""
    slide = _new_slide(prs); fh, fb = T["font_h"], T["font_b"]
    _rect(slide, 0, 0, W, H, T["light_bg"])
    _rect(slide, 0, 0, W, Inches(1.12), T["dark_bg"])
    _rect(slide, 0, Inches(1.12), W, Inches(0.048), T["accent"])
    _txt(slide, titulo, Inches(0.4), Inches(0.11), Inches(11.5), Inches(0.92),
         size=26, bold=True, color=T["dark_text"], align=PP_ALIGN.LEFT, font=fh)
    if num:
        _txt(slide, num, Inches(12.1), Inches(0.11), Inches(1.0), Inches(0.92),
             size=13, color=T["muted"], align=PP_ALIGN.RIGHT, font=fb)
    n = min(len(members), 4)
    if n==0: return slide
    cw = Inches(12.5/n)-Inches(0.14); sx = Inches(0.38); cy = Inches(1.28); ch = Inches(5.9)
    for i, mb in enumerate(members[:n]):
        nm, rl, dc = (mb[0] if len(mb)>0 else "Nombre"),(mb[1] if len(mb)>1 else ""),(mb[2] if len(mb)>2 else "")
        cx = sx + i*(cw+Inches(0.14))
        _rect(slide, cx, cy, cw, ch, T["dark_bg"])
        _rect(slide, cx, cy, cw, Inches(0.32), T["accent"])
        ar = Inches(0.72); ax = cx+cw/2-ar; ay = cy+Inches(0.5)
        _rect(slide, ax, ay, ar*2, ar*2, T["highlight"])
        _txt(slide, (nm[0].upper() if nm else "?"), ax, ay, ar*2, ar*2, size=28, bold=True, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fh)
        _txt(slide, nm, cx+Inches(0.08), ay+ar*2+Inches(0.18), cw-Inches(0.16), Inches(0.62), size=14.5, bold=True, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fh)
        _txt(slide, rl, cx+Inches(0.08), ay+ar*2+Inches(0.85), cw-Inches(0.16), Inches(0.5), size=12, color=T["muted"], align=PP_ALIGN.CENTER, font=fb)
        if dc: _txt(slide, dc, cx+Inches(0.12), ay+ar*2+Inches(1.42), cw-Inches(0.24), Inches(2.4), size=11.5, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fb)
    _notes(slide, f"{titulo}\n" + "\n".join(f"{m[0]}: {m[1]}" for m in members[:n]))
    return slide


def _L_closing(prs, titulo, cta, contacto, T, logo=None):
    """Cierre — fondo de acento, gracias, CTA, contacto."""
    slide = _new_slide(prs); fh, fb = T["font_h"], T["font_b"]
    _rect(slide, 0, 0, W, H, T["accent"])
    _rect(slide, 0, 0, Inches(0.16), H, T["dark_bg"])
    _txt(slide, "¡Gracias!", Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.0),
         size=64, bold=True, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fh)
    if cta: _txt(slide, cta, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.4), size=26, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fb)
    _rect(slide, Inches(4.2), Inches(4.98), Inches(4.9), Inches(0.038), T["dark_bg"])
    if contacto: _txt(slide, contacto, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.7), size=18, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fb)
    _txt(slide, titulo + "  ·  WEP AI", Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6), size=13, color=T["dark_text"], align=PP_ALIGN.CENTER, font=fb)
    # Logo del usuario en cierre
    if logo:
        _insert_logo(slide, logo, position='bottom_left', max_w_in=1.8, max_h_in=0.65)
    _notes(slide, f"Cierre — {cta}\n{contacto}")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# 12 PLANNERS DE DECK — devuelven lista de (tipo_slide, kwargs)
# ══════════════════════════════════════════════════════════════════════════════

def _plan_pitch(titulo, d, secciones, T, cfg=None):
    # ── v15.14.0 — Localización por país (moneda, autoridad fiscal) ──
    cfg = cfg or {}
    # ── v15.15.0 — FiscalVerifier: actualización autónoma de tasas/impuestos ──
    # v15.16.2 — verificar ANTES de leer las variables, para no congelar valores
    # viejos si el verifier detecta un cambio en esta misma generación.
    try:
        from agent.fiscal_verifier import verify_fiscal_data
        cfg = verify_fiscal_data(cfg)
    except Exception:
        pass
    sym = cfg.get("simbolo", "$")
    moneda = cfg.get("moneda", "USD")
    pais_nombre = cfg.get("pais", "")
    tax_authority = cfg.get("autoridad_fiscal", "")
    empresa = d.get("empresa","[Empresa]"); mercado = d.get("mercado","[Mercado]")
    monto   = d.get("monto","[Monto]");    fundador = d.get("fundador","[Fundador]")
    return [
        ("agenda", {"items":[(1,"El Problema",""),(2,"Solución",""),(3,"Mercado",""),(4,"Modelo de Negocio",""),(5,"Tracción",""),(6,"Equipo",""),(7,"Finanzas",""),(8,"Inversión","")], "titulo_agenda":"AGENDA"}),
        ("divider", {"numero":"01","titulo_sec":"El Problema","descripcion":"El dolor que resolvemos"}),
        ("bullets", {"titulo":"El Problema", "bullets":[f"¿Qué problema real y masivo resolvemos en {mercado}?","• El mercado actual carece de una solución eficiente, rápida y accesible","• Los usuarios pierden tiempo y dinero con alternativas",f"• Impacto: {mercado} afectados globalmente","","## Datos del problema","• Tamaño del problema: [cifra concreta]","• Costo actual para el usuario: [cifra]"]}),
        ("divider", {"numero":"02","titulo_sec":"Nuestra Solución","descripcion":"Rápida · Simple · Escalable"}),
        ("two_col", {"titulo":"La Solución — " + empresa,"left_title":"¿QUÉ HACEMOS?","left_items":[f"{empresa} resuelve esto mediante [descripción de la solución]","","## Diferenciador clave","• Ventaja 1 que ningún competidor tiene","• Ventaja 2 (tecnología, acceso, red)","• Por qué ahora es el momento"],"right_title":"CÓMO FUNCIONA","right_items":["## Paso 1","• El usuario llega con el problema","## Paso 2","• Nuestra plataforma lo resuelve","## Paso 3","• Resultado medible y recurrente"]}),
        ("divider", {"numero":"03","titulo_sec":"Mercado Objetivo","descripcion":"TAM · SAM · SOM"}),
        ("stats", {"titulo":"Tamaño del Mercado","stats":[("$0B","TAM — Total Addressable Market","Mercado total disponible"),("$0M","SAM — Serviceable","Mercado al que podemos llegar"),("$0M","SOM — Objetivo","Mercado a capturar en 3 años"),(f"{mercado}","Unidades objetivo","Clientes potenciales")]}),
        ("divider", {"numero":"04","titulo_sec":"Modelo de Negocio","descripcion":"¿Cómo generamos ingresos?"}),
        ("bullets", {"titulo":"Modelo de Negocio","bullets":["## Fuentes de Ingreso","• Suscripción mensual / anual: [precio]","• Comisión por transacción: [%]","• Enterprise / Licencia: [precio]","","## Métricas clave","• LTV (Lifetime Value): [cifra]","• CAC (Costo adquisición): [cifra]","• Margen bruto objetivo: [%]"]}),
        ("divider", {"numero":"05","titulo_sec":"Tracción","descripcion":"Lo que ya hemos logrado"}),
        ("stats", {"titulo":"Tracción — Lo que hemos logrado","stats":[("[X]","Clientes activos","En [N] meses"),("[X]%","Crecimiento MoM","Mes a mes"),("$[X]","MRR","Monthly Recurring Revenue"),("[X]","NPS","Net Promoter Score")]}),
        ("divider", {"numero":"06","titulo_sec":"El Equipo","descripcion":"Las personas detrás del proyecto"}),
        ("team", {"titulo":"Nuestro Equipo","members":[(fundador,"CEO & Co-Founder","[N] años de experiencia en [industria]. Ex-[empresa relevante]"),("[Co-Fundador]","CTO & Co-Founder","[N] años en ingeniería. Ex-[empresa]. [X] productos lanzados"),("[Director]","[Rol]","[Experiencia relevante y logros medibles"),("[Advisor]","Strategic Advisor","Ex-[empresa]. [Credential relevante]")]}),
        ("divider", {"numero":"07","titulo_sec":"Proyecciones Financieras","descripcion":"Año 1 → 3"}),
        ("bullets", {"titulo":"Proyecciones — 3 Años","bullets":["## Año 1 — Validación","• Revenue: $[X]  ·  Usuarios: [X]  ·  MRR final: $[X]","","## Año 2 — Crecimiento","• Revenue: $[X]  ·  Break-even en Q[N]","","## Año 3 — Escala","• Revenue: $[X]  ·  EBITDA positivo: [%]","","## Supuestos clave","• CAC: $[X] / Churn: [%] / Crecimiento MoM: [%]"]}),
        ("divider", {"numero":"08","titulo_sec":"La Inversión","descripcion":"Lo que buscamos"}),
        ("bullets", {"titulo":f"Buscamos — {monto} {moneda}","bullets":[f"## Ronda: {monto} {moneda}","Instrumento: [SAFE / Equity / Convertible Note]",f"Valoración pre-money: {sym}[X] {moneda}","","## Uso de los fondos","• [X]% → Producto y tecnología","• [X]% → Ventas y marketing","• [X]% → Equipo (nuevas contrataciones)","• [X]% → Operaciones","","## Milestones con esta inversión","• [Milestone 1 — fecha]","• [Milestone 2 — fecha]"]}),
    ]


def _plan_sales(titulo, d, secciones, T, cfg=None):
    # ── v15.14.0 — Localización por país (moneda, autoridad fiscal) ──
    cfg = cfg or {}
    try:
        from agent.fiscal_verifier import verify_fiscal_data
        cfg = verify_fiscal_data(cfg)
    except Exception:
        pass
    sym = cfg.get("simbolo", "$")
    moneda = cfg.get("moneda", "USD")
    pais_nombre = cfg.get("pais", "")
    tax_authority = cfg.get("autoridad_fiscal", "")
    empresa = d.get("empresa", "[Empresa]"); cliente = d.get("cliente", "[Cliente]")
    producto = d.get("producto", "[Producto]"); valor = d.get("val") or d.get("monto","[Precio]")
    return [
        ("bullets", {"titulo":f"Agenda — {titulo}","bullets":["1. Quiénes somos","2. El desafío que enfrenta "+cliente,"3. Nuestra propuesta de valor","4. Cómo funciona","5. Casos de éxito","6. Inversión","7. Próximos pasos"]}),
        ("two_col", {"titulo":"¿Quiénes Somos?","left_title":"NUESTRA EMPRESA","left_items":[f"## {empresa}",f"{empresa} — descripción de la empresa en 2-3 líneas. Misión y propuesta de valor.","","• Fundada en [año]","• [N] clientes activos","• Presencia en [N] países"],"right_title":"NUESTROS NÚMEROS","right_items":["## Tracción","• [X] clientes satisfechos","• [X]% de retención anual","• NPS promedio: [X]","","## Reconocimientos","• [Premio o certificación]","• [Partner certificado]"]}),
        ("bullets", {"titulo":f"El Desafío de {cliente}","bullets":[f"Entendemos que {cliente} hoy se enfrenta a:","","• Problema principal: [descripción concreta]","• Impacto en el negocio: [costo, tiempo perdido, riesgo]","• La solución actual no es suficiente porque [razón]","","## Consecuencias de no actuar","• [Consecuencia operacional]","• [Consecuencia financiera]","• [Riesgo estratégico]"]}),
        ("two_col", {"titulo":f"Nuestra Propuesta — {producto}","left_title":"QUÉ INCLUYE","left_items":[f"## {producto}",f"{producto} — descripción clara del producto/servicio","","• Característica principal 1","• Característica principal 2","• Característica principal 3","• Soporte / SLA incluido"],"right_title":"QUÉ GANA USTED","right_items":["## Beneficios medibles","• Ahorro de tiempo: [X] horas/mes","• Reducción de costos: [X]%","• Aumento de eficiencia: [X]%","","## Diferencial clave","• Lo que solo nosotros ofrecemos"]}),
        ("stats", {"titulo":"Resultados Comprobados con Clientes","stats":[("[X]%","Reducción de costos","Promedio en clientes similares"),("[X]x","ROI en el 1er año","Retorno sobre la inversión"),("[X]","Clientes satisfechos","En [industria]"),("[X]%","Retención a 2 años","Tasa de renovación")]}),
        ("bullets", {"titulo":"Casos de Éxito","bullets":["## Cliente 1 — [Industria]","• Desafío: [descripción]  ·  Solución: [producto]","• Resultado: [cifra concreta y medible]","","## Cliente 2 — [Industria]","• Desafío: [descripción]  ·  Solución: [producto]","• Resultado: [cifra concreta y medible]","","## Testimonio","• \"[Cita directa del cliente sobre el resultado]\" — [Nombre, Cargo]"]}),
        ("bullets", {"titulo":"Inversión","bullets":[f"## Paquete recomendado para {cliente}","","• Plan: [nombre]",f"• Inversión: {valor}  ·  Forma de pago: [mensual / anual / único]","• Incluye: [lista de lo incluido]","","## Condiciones","• Forma de pago: [mensual / anual / único]","• Implementación: [N] semanas","• SLA y soporte: [detalle]","","## Garantía","• [Política de satisfacción / prueba gratuita]"]}),
        ("bullets", {"titulo":"Próximos Pasos","bullets":["## Para avanzar juntos:","","• Paso 1 — Hoy: Confirmar interés y resolver dudas","• Paso 2 — Esta semana: Envío de propuesta formal",f"• Paso 3 — Demo personalizada de {producto} para {cliente}","• Paso 4 — Acuerdo: Firma de contrato e inicio","","## ¿Tienes preguntas?","Estamos aquí para resolverlas ahora"]}),
    ]


def _plan_report(titulo, d, secciones, T, cfg=None):
    # ── v15.14.0 — Localización por país (moneda, autoridad fiscal) ──
    cfg = cfg or {}
    try:
        from agent.fiscal_verifier import verify_fiscal_data
        cfg = verify_fiscal_data(cfg)
    except Exception:
        pass
    sym = cfg.get("simbolo", "$")
    moneda = cfg.get("moneda", "USD")
    pais_nombre = cfg.get("pais", "")
    tax_authority = cfg.get("autoridad_fiscal", "")
    # ── v15.15.0 — AccountingVerifier: actualización autónoma NIIF/NIF/IFRS ──
    try:
        from agent.accounting_verifier import verify_accounting_data
        cfg = verify_accounting_data(cfg)
    except Exception:
        pass
    empresa = d.get("empresa","[Empresa]"); periodo = d.get("periodo","[Período]"); autor = d.get("autor","[Área]")
    try:
        cbuf = _chart_bars(["KPI 1","KPI 2","KPI 3","KPI 4"],[85,92,67,78],"Cumplimiento de KPIs (%)","#185FA5","F0F4F8")
    except Exception as e:
        print(f"[WEP AI ppt] _chart_bars (report) failed: {e}")
        cbuf = None
    try:
        lbuf = _chart_lines(["Ene","Feb","Mar","Apr","May","Jun"],{"Real":[40,55,70,65,80,95],"Meta":[60,65,70,75,80,85]},"Evolución vs Meta","F0F4F8")
    except Exception as e:
        print(f"[WEP AI ppt] _chart_lines (report) failed: {e}")
        lbuf = None
    slides = [
        ("bullets", {"titulo":f"Resumen Ejecutivo — {periodo}","bullets":[f"## {empresa} — {periodo}","• Revenue / ventas — [cifra vs objetivo]","• Resultado principal 2 — [cifra y contexto]","• Resultado principal 3 — [cifra y contexto]","","## Semáforo general",f"• 🟢 Áreas de {empresa} en meta: [lista]","• 🟡 Áreas en seguimiento: [lista]","• 🔴 Áreas con alerta: [lista]"]}),
        ("stats", {"titulo":f"KPIs Principales — {periodo}","stats":[("[X]%","Cumplimiento global","vs {meta}% de meta"),("[X]","Resultado 2","Descripción"),("[X]","Resultado 3","Descripción"),("[X]%","Variación vs período anterior","YoY / MoM")]}),
    ]
    if cbuf: slides.append(("chart", {"titulo":"Desempeño por KPI","chart_buf":cbuf,"caption":"Fig 1 — Cumplimiento de indicadores clave","bullets":["## Análisis","• KPI 1: [interpretación]","• KPI 2: [interpretación]","• KPI 3: [interpretación]","","## Tendencia","• Mejora de [X]% vs período anterior"]}))
    if lbuf: slides.append(("chart", {"titulo":"Evolución vs Meta","chart_buf":lbuf,"caption":"Fig 2 — Tendencia real vs meta mensual","bullets":["## Lectura","• Superamos la meta en [N] meses","• Brecha máxima: [X] en [mes]","## Proyección","• Cierre estimado: [cifra]"]}))
    slides += [
        ("two_col", {"titulo":"Análisis por Área","left_title":"✅ LOGROS DESTACADOS","left_items":["## Área 1","• [Logro medible con cifra]","## Área 2","• [Logro medible con cifra]","## Área 3","• [Logro medible con cifra]"],"right_title":"⚠️ ÁREAS DE MEJORA","right_items":["## Desafío 1","• [Descripción + impacto]","• Acción correctiva: [qué se hará]","## Desafío 2","• [Descripción + impacto]"]}),
        ("timeline", {"titulo":"Plan de Acción — Próximos 90 Días","steps":[(1,"Acción 1","[Responsable · Fecha]"),(2,"Acción 2","[Responsable · Fecha]"),(3,"Acción 3","[Responsable · Fecha]"),(4,"Acción 4","[Responsable · Fecha]"),(5,"Revisión","[Fecha de seguimiento]")]}),
        ("bullets", {"titulo":"Conclusiones y Compromisos","bullets":["## Conclusión principal","• El período muestra [evaluación general] con [cifra clave]","","## Compromisos para el próximo período","• Compromiso 1: [acción + responsable + fecha]","• Compromiso 2: [acción + responsable + fecha]","• Compromiso 3: [acción + responsable + fecha]","","## Próxima revisión",f"• Fecha: [fecha]  ·  Responsable: {autor}"]}),
    ]
    return slides


def _plan_training(titulo, d, secciones, T, cfg=None):
    instructor = d.get("instructor","[Instructor]"); duracion = d.get("duracion","[Duración]")
    return [
        ("agenda", {"items":[(1,"Bienvenida y objetivos",""),(2,"Módulo 1 — Conceptos",""),(3,"Ejercicio práctico",""),(4,"Módulo 2 — Aplicación",""),(5,"Casos de estudio",""),(6,"Q&A y evaluación","")], "titulo_agenda":"AGENDA DEL TALLER"}),
        ("bullets", {"titulo":"Objetivos de Aprendizaje","bullets":["Al finalizar este taller el participante podrá:","","• Objetivo 1 — [habilidad concreta medible]","• Objetivo 2 — [habilidad concreta medible]","• Objetivo 3 — [habilidad concreta medible]","","## Requisitos previos","• [Conocimiento o experiencia necesaria]","",f"## Duración total: {duracion}"]}),
        ("divider", {"numero":"01","titulo_sec":"Módulo 1","descripcion":"Conceptos y fundamentos"}),
        ("bullets", {"titulo":"Módulo 1 — Conceptos Base","bullets":["## ¿Qué es [tema]?","Definición clara y concisa del concepto central","","## ¿Por qué importa?","• Razón 1 — impacto en el trabajo diario","• Razón 2 — relevancia en el contexto actual","","## Conceptos clave","• Concepto A: [definición breve]","• Concepto B: [definición breve]","• Concepto C: [definición breve]"]}),
        ("bullets", {"titulo":"Módulo 1 — ¿Cómo funciona?","bullets":["## Proceso / Metodología","• Paso 1: [descripción con acción concreta]","• Paso 2: [descripción con acción concreta]","• Paso 3: [descripción con acción concreta]","","## Ejemplo práctico","• Situación: [contexto realista]","• Aplicación: [cómo se aplica el concepto]","• Resultado: [qué se obtiene]"]}),
        ("divider", {"numero":"02","titulo_sec":"Ejercicio Práctico","descripcion":"Aplica lo aprendido"}),
        ("bullets", {"titulo":"Ejercicio — Caso Práctico","bullets":["## Instrucciones","• Tiempo: [N] minutos  ·  Equipos de: [N] personas","","## Escenario","[Descripción del caso o situación a resolver]","","## Tu tarea","• 1. Analiza la situación","• 2. Identifica el problema central","• 3. Propón una solución usando los conceptos del Módulo 1","• 4. Presenta tu solución en [N] minutos"]}),
        ("divider", {"numero":"03","titulo_sec":"Módulo 2","descripcion":"Aplicación avanzada"}),
        ("bullets", {"titulo":"Módulo 2 — Aplicación Práctica","bullets":["## Nivel avanzado","Building on Módulo 1, ahora aplicamos en contextos más complejos","","## Técnicas avanzadas","• Técnica A: [descripción + cuándo usarla]","• Técnica B: [descripción + cuándo usarla]","","## Errores comunes a evitar","• Error 1: [descripción y cómo corregirlo]","• Error 2: [descripción y cómo corregirlo]"]}),
        ("bullets", {"titulo":"Puntos Clave a Recordar","bullets":["## Las 5 ideas más importantes de hoy:","","• 1️⃣ [Concepto más importante]","• 2️⃣ [Segundo concepto más importante]","• 3️⃣ [Tercer concepto]","• 4️⃣ [Cuarto concepto]","• 5️⃣ [Quinto concepto]","","## Recursos adicionales","• [Libro / artículo / herramienta recomendada]",f"• Instructor/a: {instructor}"]}),
    ]


def _plan_company(titulo, d, secciones, T, cfg=None):
    # ── v15.14.0 — Localización por país (moneda, autoridad fiscal) ──
    cfg = cfg or {}
    try:
        from agent.corporate_verifier import verify_corporate_data
        cfg = verify_corporate_data(cfg)
    except Exception:
        pass
    sym = cfg.get("simbolo", "$")
    moneda = cfg.get("moneda", "USD")
    pais_nombre = cfg.get("pais", "")
    tax_authority = cfg.get("autoridad_fiscal", "")
    empresa = d.get("empresa","[Empresa]"); fundacion = d.get("fundacion","[Año]"); sector = d.get("sector","[Sector]")
    return [
        ("bullets", {"titulo":f"¿Quiénes somos? — {empresa}","bullets":[f"## {empresa}",f"{empresa} — descripción de la empresa. Misión, visión y propósito en 2-3 líneas.",f"","## Datos clave",f"• Fundada en: {fundacion}",f"• Fundada en {fundacion}  ·  Sector: {sector}","• Presencia: [ciudades / países]","• Empleados: [número]","","## Nuestra misión","[Declaración de misión: por qué existimos]"]}),
        ("two_col", {"titulo":"Historia y Trayectoria","left_title":"NUESTRO INICIO","left_items":["## El origen","Cómo y por qué surgió la empresa","","• [Año] — Fundación","• [Año] — Primer producto/cliente","• [Año] — Expansión","• [Año] — Hito relevante"],"right_title":"HOY","right_items":["## Dónde estamos","• [X] clientes en [N] países","• [X] productos / servicios","• [X] empleados","","## Reconocimientos","• [Certificación o premio]","• [Mención relevante]"]}),
        ("bullets", {"titulo":"Nuestros Productos y Servicios","bullets":["## Portafolio principal","","• Producto/Servicio A: [descripción breve + beneficio]","• Producto/Servicio B: [descripción breve + beneficio]","• Producto/Servicio C: [descripción breve + beneficio]","","## ¿A quién atendemos?","• Segmento 1: [descripción]","• Segmento 2: [descripción]","• Geografías: [mercados que atendemos]"]}),
        ("bullets", {"titulo":f"Datos de {empresa}","bullets":[f"## {empresa}",f"• Fundada en: {fundacion}",f"• Sector: {sector}","• Países: [ciudades / países]","• Equipo: [número] personas"]}),
        ("stats", {"titulo":"Nuestros Números","stats":[("[X]+","Clientes activos","En [N] países"),("[X]+","Años de experiencia",f"Fundada en {fundacion}"),("[X]%","Retención","Tasa de renovación anual"),("[X]","Proyectos","Completados exitosamente")]}),
        ("bullets", {"titulo":"Nuestra Propuesta de Valor","bullets":["## ¿Por qué elegirnos?","","• Diferenciador 1: [qué nos hace únicos]","• Diferenciador 2: [expertise o tecnología]","• Diferenciador 3: [modelo de atención o resultados]","","## Nuestro compromiso","• Con la calidad: [estándar específico]","• Con el cliente: [promesa concreta]","• Con la innovación: [cómo evolucionamos]"]}),
        ("bullets", {"titulo":"Alianzas y Clientes","bullets":["## Clientes que confían en nosotros","• [Cliente 1] — [industria]","• [Cliente 2] — [industria]","• [Cliente 3] — [industria]","• [Cliente 4] — [industria]","","## Alianzas estratégicas","• [Aliado 1] — [tipo de alianza]","• [Aliado 2] — [tipo de alianza]","","## Certificaciones","• [Certificación o estándar cumplido]"]}),
        ("bullets", {"titulo":"Contacto y Próximos Pasos","bullets":["## Estamos listos para trabajar contigo","","• 📧 Email: [contacto@empresa.com]","• 🌐 Web: [www.empresa.com]","• 📞 Teléfono: [número]","• 📍 Dirección: [ciudad, país]","","## ¿Cómo empezamos?","• Agenda una llamada: [enlace o contacto]","• Solicita una demo: [proceso]","• Pide tu cotización: [cómo hacerlo]"]}),
    ]


def _plan_product(titulo, d, secciones, T, cfg=None):
    # ── v15.14.0 — Localización por país (moneda, autoridad fiscal) ──
    cfg = cfg or {}
    try:
        from agent.fiscal_verifier import verify_fiscal_data
        cfg = verify_fiscal_data(cfg)
    except Exception:
        pass
    sym = cfg.get("simbolo", "$")
    moneda = cfg.get("moneda", "USD")
    pais_nombre = cfg.get("pais", "")
    tax_authority = cfg.get("autoridad_fiscal", "")
    producto = d.get("producto","[Producto]"); empresa = d.get("empresa","[Empresa]")
    return [
        ("bullets", {"titulo":f"Bienvenidos a {producto} — por {empresa}","bullets":[f"## {producto} — por {empresa}",f"La solución definitiva para [el problema que resuelve {producto}]","","## El problema que resolvemos","[Descripción del dolor que experimenta tu usuario objetivo]","","## Nuestra solución",f"{producto} — [cómo resuelves el problema de forma única]","","## ¿Por qué ahora?","[Razón de timing — tendencia de mercado, cambio tecnológico]"]}),
        ("bullets", {"titulo":f"¿Qué hace {producto}?","bullets":["## Funcionalidades principales","","• Feature 1: [descripción clara y beneficio]","• Feature 2: [descripción clara y beneficio]","• Feature 3: [descripción clara y beneficio]","• Feature 4: [descripción clara y beneficio]","","## Integrations","• Conecta con: [herramienta], [herramienta], [herramienta]","• API disponible: [REST / GraphQL / Webhooks]"]}),
        ("timeline", {"titulo":f"¿Cómo funciona {producto}?","steps":[(1,"Onboarding","Registro en [X] minutos"),(2,"Configuración","Setup guiado paso a paso"),(3,"Integración","Conecta tus herramientas"),(4,"Uso","Empieza a generar valor"),(5,"Resultados","Mide el impacto")]}),
        ("two_col", {"titulo":"Para quién es","left_title":"USUARIO IDEAL","left_items":["## Perfil","• [Cargo o tipo de usuario]","• [Tamaño de empresa]","• [Industria]","","## Necesidad","• [Pain point 1]","• [Pain point 2]","• [Resultado que busca]"],"right_title":"CASOS DE USO","right_items":["## Uso 1","• [Descripción de cómo lo usa]","## Uso 2","• [Descripción de cómo lo usa]","## Uso 3","• [Descripción de cómo lo usa]"]}),
        ("stats", {"titulo":"Por qué los usuarios nos eligen","stats":[("[X]min","Setup","Hasta el primer resultado"),("[X]%","Usuarios activos","Retención a 6 meses"),("[X]%","NPS","Net Promoter Score"),("[X]x","ROI","Promedio en 90 días")]}),
        ("bullets", {"titulo":"Planes y Precios","bullets":["## Free / Starter","• [X] [unidad] gratis  ·  Sin tarjeta de crédito","• Incluye: [funcionalidades básicas]","","## Pro — $[X]/mes","• [X] [unidad]  ·  Todas las features","• Soporte: [nivel]","","## Enterprise — Personalizado","• Unlimited · SLA garantizado · Implementación dedicada","","## ¿Cuál es el correcto para ti?","Agenda una llamada y te ayudamos a elegir"]}),
    ]


def _plan_strategy(titulo, d, secciones, T, cfg=None):
    # ── v15.14.0 — Localización por país (moneda, autoridad fiscal) ──
    cfg = cfg or {}
    try:
        from agent.fiscal_verifier import verify_fiscal_data
        cfg = verify_fiscal_data(cfg)
    except Exception:
        pass
    sym = cfg.get("simbolo", "$")
    moneda = cfg.get("moneda", "USD")
    pais_nombre = cfg.get("pais", "")
    tax_authority = cfg.get("autoridad_fiscal", "")
    empresa = d.get("empresa","[Empresa]"); horizonte = d.get("horizonte","3 años")
    monto = d.get("monto", "[Inversión estimada]")
    return [
        ("bullets", {"titulo":"Contexto Estratégico","bullets":["## Punto de partida",f"• {empresa} hoy: [descripción del estado actual]","• Qué ha cambiado: [factores del entorno]","","## El mandato","[Descripción del objetivo estratégico en 1-2 líneas]","",f"## Horizonte de planificación: {horizonte}"]}),
        ("two_col", {"titulo":"Análisis FODA","left_title":"FORTALEZAS / OPORTUNIDADES","left_items":["## Fortalezas internas","• [F1] — [impacto]","• [F2] — [impacto]","• [F3] — [impacto]","","## Oportunidades externas","• [O1] — [potencial]","• [O2] — [potencial]"],"right_title":"DEBILIDADES / AMENAZAS","right_items":["## Debilidades a superar","• [D1] — [plan de mejora]","• [D2] — [plan de mejora]","","## Amenazas a mitigar","• [A1] — [plan de mitigación]","• [A2] — [plan de mitigación]"]}),
        ("bullets", {"titulo":"Pilares Estratégicos","bullets":["## Los [N] ejes de nuestra estrategia:","","## 1. [Pilar 1 — ej: Crecimiento]","• [Descripción de la apuesta estratégica]","• Meta: [indicador medible]","","## 2. [Pilar 2 — ej: Eficiencia]","• [Descripción de la apuesta estratégica]","• Meta: [indicador medible]","","## 3. [Pilar 3 — ej: Innovación]","• [Descripción de la apuesta estratégica]","• Meta: [indicador medible]"]}),
        ("stats", {"titulo":"Metas Estratégicas — " + horizonte,"stats":[("[X]","Meta 1","[Unidad y descripción]"),("[X]%","Meta 2","[Descripción]"),("[X]","Meta 3","[Descripción]"),("[X]x","Meta 4","[Descripción]")]}),
        ("timeline", {"titulo":"Hoja de Ruta — Roadmap","steps":[(1,"Año 1\nFundamentos","Bases y quick wins"),(2,"Año 2\nCrecimiento","Escala y optimización"),(3,"Año 3\nLiderazgo","Consolidación y expansión")]}),
        ("bullets", {"titulo":"Recursos e Inversión Requerida","bullets":["## Capital humano","• Nuevas contrataciones: [perfiles y número]","• Capacitación: [inversión y programa]","","## Capital financiero",f"• Inversión total: {monto} (estimado)","• Fuentes: [interna / externa / mixta]","","## Tecnología","• [Herramienta/sistema] — [costo estimado]","","## KPIs de seguimiento trimestral","• [KPI 1], [KPI 2], [KPI 3]"]}),
    ]


def _plan_board(titulo, d, secciones, T, cfg=None):
    # ── v15.14.0 — Localización por país (moneda, autoridad fiscal) ──
    cfg = cfg or {}
    try:
        from agent.fiscal_verifier import verify_fiscal_data
        cfg = verify_fiscal_data(cfg)
    except Exception:
        pass
    sym = cfg.get("simbolo", "$")
    moneda = cfg.get("moneda", "USD")
    pais_nombre = cfg.get("pais", "")
    tax_authority = cfg.get("autoridad_fiscal", "")
    # ── v15.15.0 — AccountingVerifier: actualización autónoma NIIF/NIF/IFRS ──
    try:
        from agent.accounting_verifier import verify_accounting_data
        cfg = verify_accounting_data(cfg)
    except Exception:
        pass
    empresa = d.get("empresa","[Empresa]"); periodo = d.get("periodo","[Período]")
    try:
        cbuf = _chart_bars(["Q1","Q2","Q3","Q4"],[82,88,79,95],"Revenue vs Plan (%)","#185FA5","FAFAFA")
    except Exception as e:
        print(f"[WEP AI ppt] _chart_bars (board) failed: {e}")
        cbuf = None
    slides = [
        ("bullets", {"titulo":f"Estado General {empresa}  ·  {periodo}","bullets":[f"## {empresa} — {periodo}","• Resultado 1: [cifra + evaluación]","• Resultado 2: [cifra + evaluación]","• Resultado 3: [cifra + evaluación]","","## Semáforo ejecutivo","• 🟢 En plan: [áreas]","• 🟡 Monitoreo: [áreas]","• 🔴 Requiere atención: [áreas]"]}),
        ("stats", {"titulo":"Resultados Financieros","stats":[("$[X]M","Revenue","vs $[X]M plan"),("[X]%","Margen bruto","vs [X]% período anterior"),("$[X]M","EBITDA","[X]% del revenue"),("[X]%","Burn rate","Runway: [N] meses")]}),
    ]
    if cbuf: slides.append(("chart", {"titulo":"Revenue vs Plan — " + periodo,"chart_buf":cbuf,"caption":"Fig 1 — Desempeño trimestral vs presupuesto","bullets":["## Análisis","• Q mejor: [Q] — [razón]","• Q desafío: [Q] — [razón]","• Tendencia: [positiva/negativa]","## Proyección cierre","• [Cifra y escenarios]"]}))
    slides += [
        ("two_col", {"titulo":"Estrategia — Avances y Desvíos","left_title":"✅ EN PLAN","left_items":["## Iniciativas cumplidas","• [Iniciativa 1] — [resultado]","• [Iniciativa 2] — [resultado]","","## Hitos alcanzados","• [Hito 1] — [fecha]","• [Hito 2] — [fecha]"],"right_title":"⚠️ DESVÍOS Y PLAN","right_items":["## Iniciativas con desvío","• [Iniciativa] — [razón del desvío]","• Plan correctivo: [acción]","","## Riesgos activos","• [Riesgo 1] — [probabilidad / impacto]","• [Riesgo 2] — [plan de mitigación]"]}),
        ("bullets", {"titulo":"Decisiones Requeridas a la Junta","bullets":["## Puntos de aprobación","","• 1. [Decisión 1 — descripción + recomendación de mgmt]","• 2. [Decisión 2 — descripción + recomendación de mgmt]","• 3. [Decisión 3 — descripción + opciones a evaluar]","","## Siguientes pasos",f"• Próxima junta: [fecha]  ·  Empresa: {empresa}  ·  Período: {periodo}","• Reportes entre juntas: [frecuencia y responsable]"]}),
    ]
    return slides


def _plan_investor_update(titulo, d, secciones, T, cfg=None):
    # ── v15.14.0 — Localización por país (moneda, autoridad fiscal) ──
    cfg = cfg or {}
    try:
        from agent.fiscal_verifier import verify_fiscal_data
        cfg = verify_fiscal_data(cfg)
    except Exception:
        pass
    sym = cfg.get("simbolo", "$")
    moneda = cfg.get("moneda", "USD")
    pais_nombre = cfg.get("pais", "")
    tax_authority = cfg.get("autoridad_fiscal", "")
    # ── v15.15.0 — AccountingVerifier: actualización autónoma NIIF/NIF/IFRS ──
    try:
        from agent.accounting_verifier import verify_accounting_data
        cfg = verify_accounting_data(cfg)
    except Exception:
        pass
    empresa = d.get("empresa","[Empresa]"); mes = d.get("periodo","[Mes/Período]")
    try:
        cbuf = _chart_lines(["Ene","Feb","Mar","Apr","May","Jun"],{"MRR":[20,28,35,42,51,60],"Meta":[25,30,35,40,45,50]},"MRR vs Meta ($K)","F0F4F8")
    except Exception as e:
        print(f"[WEP AI ppt] _chart_lines (investor) failed: {e}")
        cbuf = None
    slides = [
        ("stats", {"titulo":f"Resumen ejecutivo {empresa}  ·  {mes}","stats":[("$[X]K","MRR","Monthly Recurring Revenue"),("[X]%","MoM Growth","Crecimiento mensual"),("[X]","Clientes activos","Total acumulado"),("[X] meses","Runway","Con burn actual")]}),
    ]
    if cbuf: slides.append(("chart", {"titulo":"MRR — Evolución","chart_buf":cbuf,"caption":"Crecimiento mensual vs meta","bullets":["## Highlights","• Mejor mes: [mes] — $[X]K","• Cierre proyectado: $[X]K","## Clientes","• Nuevos este mes: [N]","• Churn: [N] ([X]%)"]}))
    slides += [
        ("two_col", {"titulo":"Highlights & Lowlights","left_title":"🟢 HIGHLIGHTS","left_items":["## Lo que salió bien","• [Logro 1 con cifra]","• [Logro 2 con cifra]","• [Logro 3 con cifra]","","## Hito del mes","• [Hito más relevante]"],"right_title":"🔴 LOWLIGHTS","right_items":["## Desafíos del mes","• [Desafío 1 — qué haremos]","• [Desafío 2 — qué haremos]","","## Aprendizajes","• [Qué aprendimos]"]}),
        ("bullets", {"titulo":"Métricas Detalladas","bullets":["## Revenue","• MRR: $[X]K (+[X]% MoM)  ·  ARR: $[X]K","• New MRR: $[X]K  ·  Expansion: $[X]K  ·  Churn: $[X]K","","## Clientes","• Total activos: [X]  ·  Nuevos: [X]  ·  Churned: [X]","• NPS: [X]  ·  CSAT: [X]%","","## Unit Economics","• CAC: $[X]  ·  LTV: $[X]  ·  LTV/CAC: [X]x","• Payback period: [N] meses"]}),
        ("bullets", {"titulo":"Próximo Mes — Plan","bullets":["## Prioridades","• 1. [Prioridad 1 — OKR o iniciativa]","• 2. [Prioridad 2 — OKR o iniciativa]","• 3. [Prioridad 3 — OKR o iniciativa]","","## Asks para los inversores","• [Ask 1 — intro, consejo, recurso]","• [Ask 2 — intro, consejo, recurso]","","## Próxima actualización","• Fecha: [fecha]  ·  Formato: [call / email]"]}),
    ]
    return slides


def _plan_case_study(titulo, d, secciones, T, cfg=None):
    cliente = d.get("cliente","[Cliente]"); producto = d.get("producto","[Producto/Servicio]")
    empresa = d.get("empresa", "[Empresa]")
    return [
        ("bullets", {"titulo":f"Caso de Éxito — {cliente}","bullets":[f"## {cliente}",f"{cliente} es una empresa/organización dedicada a [industria]. [Contexto relevante: tamaño, presencia, situación antes de trabajar juntos]","","## El contexto","• Industria: [sector]","• Empleados: [número]","• Presencia: [mercados]","","## Por qué nos eligieron","• [Razón principal de la selección]"]}),
        ("two_col", {"titulo":"El Desafío","left_title":"SITUACIÓN INICIAL","left_items":["## Antes de trabajar juntos","• [Problema 1 — descripción concreta]","• [Problema 2 — impacto medible]","• [Problema 3 — consecuencia operacional]","","## El costo del problema","• [Cifra económica / tiempo perdido]"],"right_title":"LO QUE BUSCABAN","right_items":["## Objetivos del cliente","• Objetivo 1: [resultado esperado]","• Objetivo 2: [resultado esperado]","• Objetivo 3: [resultado esperado]","","## Timeline","• Implementación en: [N semanas/meses]"]}),
        ("timeline", {"titulo":"La Implementación","steps":[(1,"Diagnóstico","[N] días — mapeo del estado actual"),(2,"Diseño","Configuración de la solución"),(3,"Piloto","Prueba con equipo reducido"),(4,"Rollout","Despliegue completo"),(5,"Optimización","Ajustes basados en datos")]}),
        ("stats", {"titulo":f"Los Resultados — {cliente}","stats":[("[X]%","Reducción de costos","En el área [X]"),("[X]x","ROI en [N] meses","Retorno sobre la inversión"),("[X]h","Ahorro semanal","Por empleado / área"),("[X]%","Satisfacción","NPS del equipo")]}),
        ("quote", {"quote":"[Insertar cita directa del cliente sobre el resultado y la experiencia con la solución]","author":"[Nombre del contacto]","role":"[Cargo] — "+cliente}),
        ("bullets", {"titulo":f"¿Quieres resultados como los de {cliente}?","bullets":["## Lo que logramos juntos","• [Resultado más impactante con cifra]","• [Segundo resultado]","• [Tercer resultado]","","## Próximos pasos para ti",f"• 1. Agenda una llamada sobre {producto} (30 min, gratis)","• 2. Recibe una propuesta personalizada","• 3. Empieza en [N] semanas","","## Contacto",f"• Contacto: [email]  ·  [teléfono]\n• {empresa} — [web]"]}),
    ]


def _plan_webinar(titulo, d, secciones, T, cfg=None):
    speaker = d.get("speaker","[Speaker]"); empresa = d.get("empresa","[Empresa]")
    return [
        ("bullets", {"titulo":"Bienvenidos — Agenda de hoy","bullets":[f"## {titulo}","","• ⏱ Duración: [N] minutos + Q&A abierto",f"• 🎙 Speaker: {speaker} — {empresa}","","## Lo que veremos hoy:","• 1. [Tema 1] — [X] min","• 2. [Tema 2] — [X] min","• 3. [Tema 3] — [X] min","• 4. Q&A — [X] min","","## Para participar:","• Usa el chat para preguntas"]}),
        ("bullets", {"titulo":f"¿Quién soy yo? — {speaker}","bullets":[f"## {speaker}",f"[Cargo] en {empresa}","","• [Credencial o experiencia 1]","• [Credencial o experiencia 2]","• [Logro o reconocimiento relevante]","","## Por qué este tema","[Razón personal o profesional para hablar de esto]"]}),
        ("two_col", {"titulo":"El panorama actual","left_title":"EL PROBLEMA","left_items":["## Lo que está pasando","• [Tendencia 1 — datos]","• [Tendencia 2 — impacto]","• [Tendencia 3 — urgencia]","","## Por qué importa ahora","• [Razón de timing]"],"right_title":"LA OPORTUNIDAD","right_items":["## Lo que es posible","• [Oportunidad 1]","• [Oportunidad 2]","","## Quienes ya lo están haciendo","• [Ejemplo de empresa/persona exitosa]"]}),
        ("bullets", {"titulo":"Punto 1 — [Primer tema clave]","bullets":["## La idea central","[Explicación del concepto en 2-3 líneas]","","## ¿Cómo aplicarlo?","• Paso 1: [acción concreta]","• Paso 2: [acción concreta]","• Paso 3: [acción concreta]","","## Ejemplo real","[Caso concreto con resultado medible]"]}),
        ("bullets", {"titulo":"Punto 2 — [Segundo tema clave]","bullets":["## La idea central","[Explicación del concepto]","","## Framework / Modelo","• [Elemento 1] — [descripción]","• [Elemento 2] — [descripción]","• [Elemento 3] — [descripción]","","## Aplicación práctica","• [Cómo empezar esta semana]"]}),
        ("bullets", {"titulo":"Punto 3 — [Tercer tema clave]","bullets":["## La idea central","[Explicación del concepto]","","## Los 3 pasos clave","• Paso 1: [descripción]","• Paso 2: [descripción]","• Paso 3: [descripción]","","## Herramientas recomendadas","• [Herramienta 1] — [para qué sirve]","• [Herramienta 2] — [para qué sirve]"]}),
        ("bullets", {"titulo":"Resumen y Acción Inmediata","bullets":["## Las 3 ideas más importantes de hoy:","","• 1️⃣ [Idea 1 — en una línea]","• 2️⃣ [Idea 2 — en una línea]","• 3️⃣ [Idea 3 — en una línea]","","## Tu acción para esta semana:","[Acción concreta y específica que pueden hacer en los próximos 7 días]","","## Recursos gratuitos","• [Link a recurso 1]","• [Link a recurso 2]"]}),
    ]


def _plan_generic(titulo, d, secciones, T, cfg=None):
    temas = secciones if secciones else ["Introducción","Contexto","Desarrollo","Análisis","Conclusiones y Pasos"]
    slides = []
    for i, tema in enumerate(temas):
        slides.append(("bullets", {"titulo":tema,"bullets":[f"## {tema}","","• Punto principal 1","• Punto principal 2","• Punto principal 3","","## Detalle","• Sub-punto A","• Sub-punto B"]}))
    return slides


# ══════════════════════════════════════════════════════════════════════════════
# ROUTER PRINCIPAL
# ══════════════════════════════════════════════════════════════════════════════

def generate_powerpoint(gen_data: dict) -> str:
    _ensure_dir()
    titulo     = gen_data.get("titulo", "Presentación")
    instruc    = gen_data.get("instrucciones", "")
    datos      = gen_data.get("datos", {})
    secciones  = datos.get("secciones", [])
    estilo_str = datos.get("estilo", "profesional")
    detalles   = datos.get("detalles_especificos", "")
    fecha_str  = datetime.now().strftime("%d de %B de %Y")

    # v14.4 — setear idioma para toda la generación (thread-local)
    _pp_set_lang(datos.get("idioma") or gen_data.get("idioma") or "es")

    # User data
    d = {
        "empresa":   _d(datos,"empresa") or _d(gen_data,"parte_a_nombre") or _d(gen_data,"patron_nombre","[Empresa]"),
        "cliente":   _d(datos,"cliente") or _d(gen_data,"parte_b_nombre","[Cliente]"),
        "producto":  _d(datos,"producto") or _d(gen_data,"objeto") or titulo,
        "autor":     _d(datos,"autor")   or _d(gen_data,"patron_nombre",""),
        "speaker":   _d(datos,"speaker") or _d(gen_data,"trabajador_nombre",""),
        "monto":     _d(datos,"monto")   or _d(gen_data,"monto","[Monto]"),
        "val":       _d(datos,"val")   or _d(datos,"valor") or _d(gen_data,"valor",""),
        "periodo":   _d(datos,"periodo") or detalles[:40] or "[Período]",
        "fundador":  _d(datos,"fundador") or _d(gen_data,"patron_nombre","[Fundador]"),
        "fundacion": _d(datos,"fundacion","[Año]"),
        "sector":    _d(datos,"sector",   "[Sector]"),
        "horizonte": _d(datos,"horizonte","3 años"),
        "duracion":  _d(datos,"duracion", "[Duración]"),
        "instructor":_d(datos,"instructor") or _d(gen_data,"patron_nombre","[Instructor]"),
        "mercado":   _d(datos,"mercado",  "[Mercado objetivo]"),
        "inicio":    _d(datos, "inicio",    ""),
        "nuevo":     _d(datos, "nuevo",     ""),
        "anios":     _d(datos, "anios",     "[N]"),
        "tutor":     _d(datos, "tutor",     ""),
        "pm":        _d(datos, "pm",        "") or _d(datos, "autor", "[PM]"),
        "proyecto":  (datos.get("proyecto") or "") or _d(datos, "producto", titulo),
    }

    logo_path = gen_data.get("logo") or datos.get("logo") or None

    T   = _resolve_theme(estilo_str)
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    full = f"{instruc} {detalles} {titulo}".lower()

    # ── v14.2: TEMPLATE_ID EXPLÍCITO (LLM-driven) ─────────────────────────────
    # Si Claude emitió un template_id válido, lo usamos directamente.
    template_id = (gen_data.get("template_id") or "").strip().lower()
    if template_id and template_id in PPT_TEMPLATES:
        planner = PPT_TEMPLATES[template_id]
    # Detect deck type (fallback por keywords)
    elif any(w in full for w in ["pitch","startup","inversores","ronda","seed","series a","capital"]):
        planner = _plan_pitch
    elif any(w in full for w in ["ventas","comercial","propuesta comercial","sales"]):
        planner = _plan_sales
    elif any(w in full for w in ["reporte","informe","resultado","kpi","gerencial","board report","management report"]):
        planner = _plan_report
    elif any(w in full for w in ["capacitación","capacitacion","curso","formación","taller","entrenamiento","training","workshop"]):
        planner = _plan_training
    elif any(w in full for w in ["perfil de empresa","perfil empresa","company profile","presentación institucional","institucional","about us"]):
        planner = _plan_company
    elif any(w in full for w in ["demo","product demo","demo de producto","producto saas","lanzamiento de producto"]):
        planner = _plan_product
    elif any(w in full for w in ["estrategia","plan estratégico","strategic plan","roadmap","hoja de ruta","foda"]):
        planner = _plan_strategy
    elif any(w in full for w in ["junta directiva","board deck","directorio","consejo board","board meeting","junta de socios"]):
        planner = _plan_board
    elif any(w in full for w in ["investor update","update inversores","reporte inversores","monthly update","quarterly update"]):
        planner = _plan_investor_update
    elif any(w in full for w in ["caso de éxito","caso de exito","case study","success story","cliente exitoso"]):
        planner = _plan_case_study
    elif any(w in full for w in ["webinar","conferencia","charla","keynote","ponencia","presentación pública"]):
        planner = _plan_webinar
    elif any(w in full for w in ["marca personal","personal brand","personal branding",
                                   "quién soy yo","quien soy yo","mi historia profesional",
                                   "presentación de conferencia","speaker deck","about me"]):
        planner = _plan_personal_brand

    elif any(w in full for w in ["presentación de empleo","presentacion de empleo",
                                  "entrevista de trabajo","interview presentation",
                                  "preséntate en","presentate en","proceso de selección",
                                  "selección de personal","presentación laboral"]):
        planner = _plan_job_presentation

    elif any(w in full for w in ["presentación académica","presentacion academica",
                                  "trabajo final de grado","tfg","tesis","proyecto de clase",
                                  "defensa de tesis","proyecto académico","trabajo universitario",
                                  "investigación académica","seminario académico"]):
        planner = _plan_academic

    elif any(w in full for w in ["onboarding de cliente","client onboarding","bienvenida cliente",
                                  "kick-of","kickof","inicio de proyecto","arranque de proyecto",
                                  "primer día con cliente","welcome deck"]):
        planner = _plan_client_onboarding

    elif any(w in full for w in ["all hands","all-hands","town hall","townhall",
                                  "reunión general","reunion general","asamblea de equipo",
                                  "comunicado a todo el equipo","reunión de empresa"]):
        planner = _plan_all_hands

    elif any(w in full for w in ["upsell","renovación","renovacion","expansion de cuenta",
                                  "expansión de cuenta","propuesta de renovación","cross sell",
                                  "ampliación de contrato","upgrade","retención de cliente"]):
        planner = _plan_upsell

    else:
        planner = _plan_generic

    # ── v15.13.0 — Obtener cfg de país para localización ──
    # Igual que Word y Excel: prefer datos.pais > gen_data.pais > GENERIC
    try:
        from office.legal_config import LEGAL_CONFIG, _get_legal_config
        pais_input = datos.get("pais") or gen_data.get("pais") or ""
        cfg_ppt = None
        code_ppt = None
        if pais_input:
            # Buscar por código o por nombre
            if pais_input.upper() in LEGAL_CONFIG:
                code_ppt = pais_input.upper()
                cfg_ppt  = dict(LEGAL_CONFIG[code_ppt])
            else:
                for code, entry in LEGAL_CONFIG.items():
                    if entry.get("pais", "").lower() == pais_input.lower():
                        code_ppt = code
                        cfg_ppt  = dict(entry)
                        break
        if not cfg_ppt:
            # Sin país → GENERIC (ya aplica overlay internamente)
            cfg_ppt = _get_legal_config("GENERIC")
        else:
            # v15.16.2 — aplicar el overlay persistente verificado por IA, igual
            # que hacen Word y Excel. Antes PPT leía LEGAL_CONFIG directo y las
            # actualizaciones autónomas no llegaban a las presentaciones.
            try:
                from agent.legal_overrides import apply_overlay
                cfg_ppt = apply_overlay(code_ppt, cfg_ppt)
            except Exception:
                pass
    except Exception:
        cfg_ppt = {}

    # Llamar al planner con cfg (todos los plans v15.13.0+ aceptan cfg=None)
    import inspect
    sig = inspect.signature(planner)
    if "cfg" in sig.parameters:
        slides_plan = planner(titulo, d, secciones, T, cfg=cfg_ppt)
    else:
        slides_plan = planner(titulo, d, secciones, T)

    subtitulo = detalles[:100] if detalles else instruc[:100]
    _L_cover(prs, titulo, subtitulo, d["autor"] or d["empresa"], fecha_str, T, logo=logo_path)

    for slide_type, kwargs in slides_plan:
        if slide_type == "bullets":
            num = f"{slides_plan.index((slide_type,kwargs))+1}/{len(slides_plan)}"
            _L_bullets(prs, kwargs["titulo"], kwargs.get("bullets",[]), T, num=num)
        elif slide_type == "two_col":
            _L_two_col(prs, kwargs["titulo"], kwargs.get("left_title",""), kwargs.get("left_items",[]), kwargs.get("right_title",""), kwargs.get("right_items",[]), T)
        elif slide_type == "stats":
            _L_stats(prs, kwargs["titulo"], kwargs.get("stats",[]), T)
        elif slide_type == "timeline":
            _L_timeline(prs, kwargs["titulo"], kwargs.get("steps",[]), T)
        elif slide_type == "divider":
            _L_divider(prs, kwargs.get("numero",""), kwargs.get("titulo_sec",""), kwargs.get("descripcion",""), T)
        elif slide_type == "quote":
            _L_quote(prs, kwargs.get("quote",""), kwargs.get("author",""), kwargs.get("role",""), T)
        elif slide_type == "agenda":
            _L_agenda(prs, kwargs.get("items",[]), kwargs.get("titulo_agenda","AGENDA"), T)
        elif slide_type == "chart":
            _L_chart(prs, kwargs["titulo"], kwargs.get("chart_buf"), kwargs.get("caption",""), kwargs.get("bullets",[]), T)
        elif slide_type == "team":
            _L_team(prs, kwargs["titulo"], kwargs.get("members",[]), T)

    cta     = datos.get("cta", "Quedamos a su disposición para cualquier consulta.")
    contact = datos.get("contacto", d["empresa"])
    _L_closing(prs, titulo, cta, contact, T, logo=logo_path)

    # ── v14.4 — BILINGUAL POST-PROCESSOR ──────────────────────────────────────
    # Traduce todos los runs de texto a inglés cuando el idioma del usuario es EN.
    # Es defensivo: cualquier fallo aquí no debe romper la generación de la
    # presentación; en ese caso se queda en español.
    try:
        _translate_presentation(prs)
    except Exception:
        pass

    filename = f"{_safe_filename(titulo)}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    return filepath


def generate_pdf_ppt(gen_data: dict) -> str:
    """Genera una presentación .pptx y la convierte a PDF con LibreOffice headless.
    Retorna la ruta del .pdf (o del .pptx si LibreOffice no está disponible).
    """
    import shutil

    # 1. Generar el .pptx
    pptx_path = generate_powerpoint(gen_data)

    # 2. Buscar LibreOffice
    lo = shutil.which('libreoffice') or shutil.which('soffice')
    if not lo:
        print("[WEP AI PPT] LibreOffice no disponible — entregando .pptx")
        return pptx_path

    # 3. Convertir con LibreOffice headless
    out_dir = os.path.dirname(pptx_path)
    try:
        result = subprocess.run(
            [lo, '--headless', '--convert-to', 'pdf',
             '--outdir', out_dir, pptx_path],
            capture_output=True, text=True, timeout=45
        )
        pdf_path = pptx_path.replace('.pptx', '.pdf')
        if os.path.exists(pdf_path):
            return pdf_path
        # Algunos builds de LO escriben al cwd
        local_pdf = os.path.basename(pptx_path).replace('.pptx', '.pdf')
        if os.path.exists(local_pdf):
            import shutil as sh
            sh.move(local_pdf, pdf_path)
            return pdf_path
    except Exception as e:
        print(f"[WEP AI PPT] PDF conversion error: {e}")

    return pptx_path   # Fallback a .pptx


def open_in_powerpoint(filepath: str):
    """Open the file in Microsoft PowerPoint on macOS.

    v14.8 — Usa `open -a` en lugar de AppleScript interpolado para evitar
    inyección por contenido del path. `open` recibe el path como argv y no
    lo evalúa como código.
    """
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"PowerPoint file not found: {filepath}")
    subprocess.run(["open", "-a", "Microsoft PowerPoint", filepath], check=False)


def apply_correction(filepath: str, correction: str, gen_data: dict) -> str:
    gen_data["instrucciones"] = gen_data.get("instrucciones","") + f"\n\nCORRECCIÓN: {correction}"
    return generate_powerpoint(gen_data)


# ══════════════════════════════════════════════════════════════════════════════
# GAPS v12 — 6 NUEVOS DECK TYPES
# ══════════════════════════════════════════════════════════════════════════════

def _plan_personal_brand(titulo, d, secciones, T, cfg=None):
    """Marca personal — para conferencias, podcasts, colaboraciones y visibilidad.
    No es un portafolio ni una propuesta: es narrativa personal."""
    nombre   = d.get("speaker") or d.get("autor") or d.get("empresa") or "[Tu Nombre]"
    rol      = d.get("producto") or "[Tu especialidad / rol]"
    logro    = d.get("monto")   or "[Tu logro más impactante]"
    return [
        ("quote", {
            "quote":  f"La frase que mejor define mi filosofía como {rol}",
            "author": nombre,
            "role":   rol,
        }),
        ("bullets", {
            "titulo": f"¿Quién es {nombre}?",
            "bullets": [
                f"## {nombre}",
                rol,
                "",
                "• [Experiencia 1 — con cifra o resultado concreto]",
                "• [Experiencia 2 — con cifra o resultado concreto]",
                "• [Experiencia 3 — hito o reconocimiento]",
                "",
                "## Lo que me hace diferente",
                f"• [Diferenciador 1 de {nombre} — algo que pocas personas en {rol} hacen]",
                "• [Diferenciador 2 — perspectiva única o metodología propia]",
            ],
        }),
        ("two_col", {
            "titulo": "Mi Historia",
            "left_title": "DE DÓNDE VENGO",
            "left_items": [
                "## El inicio",
                "• [Punto de partida — de dónde venías antes de llegar aquí]",
                "• [El momento que cambió tu trayectoria]",
                "• [Desafío que superaste y qué aprendiste]",
                "",
                "## La transformación",
                "• [Cómo evolucionaste de A a B]",
            ],
            "right_title": "A DÓNDE VOY",
            "right_items": [
                "## Mi misión",
                "• [Por qué haces lo que haces — propósito detrás del trabajo]",
                "",
                "## Mi visión",
                "• [Qué quieres lograr en los próximos 3-5 años]",
                "",
                "## Cómo puedo ayudarte",
                "• [Problema específico que resuelves para tu audiencia]",
            ],
        }),
        ("stats", {
            "titulo": "Mi Impacto en Números",
            "stats": [
                (logro,        "Logro principal",     "El resultado más representativo"),
                ("[X]+",       "Clientes/Proyectos",  "En [N] años de trayectoria"),
                ("[X]K",       "Comunidad",           "Seguidores / suscriptores"),
                ("[X]",        "Países",              "Donde he trabajado o presentado"),
            ],
        }),
        ("bullets", {
            "titulo": "Mis Áreas de Expertise",
            "bullets": [
                "## Lo que domino",
                "",
                "## 01 — [Área principal]",
                "• [Descripción de tu expertise + evidencia concreta]",
                "",
                "## 02 — [Segunda área]",
                "• [Descripción + metodología o herramienta propia]",
                "",
                "## 03 — [Tercera área]",
                "• [Descripción + resultado que produces con esto]",
            ],
        }),
        ("quote", {
            "quote":  "[Cita de un cliente, colega o líder de tu industria que valide tu trabajo]",
            "author": "[Nombre del que da el testimonio]",
            "role":   "[Cargo y empresa] — cliente / colaborador",
        }),
        ("bullets", {
            "titulo": "¿Cómo podemos colaborar?",
            "bullets": [
                "## Lo que ofrezco",
                "",
                "• 🎤 Conferencias y keynotes — [tema y duración]",
                "• 🧑‍🏫 Talleres y formaciones — [formato y audiencia]",
                "• 🤝 Consultoría / Mentoring — [cómo funciona]",
                "• 📝 Contenido / Colaboraciones — [tipo de colaboración]",
                "",
                "## Contáctame",
                f"• [email]  ·  [LinkedIn]  ·  [Web]",
                f"• Ubicación: [Ciudad / Remoto global]",
            ],
        }),
    ]


def _plan_job_presentation(titulo, d, secciones, T, cfg=None):
    """Presentación de empleo / entrevista — 'Preséntate en 5 minutos'.
    Solicitada en procesos de selección modernos, muy diferente a un CV."""
    nombre   = d.get("speaker") or d.get("autor") or "[Tu Nombre]"
    puesto   = d.get("producto") or d.get("rol") or "[Puesto al que aplicas]"
    empresa  = d.get("cliente") or d.get("empresa") or "[Nombre de la empresa]"
    anios    = d.get("anios") or "[N]"
    return [
        ("bullets", {
            "titulo": f"Hola, soy {nombre}",
            "bullets": [
                f"## {nombre}",
                f"{puesto}",
                "",
                f"## ¿Por qué {nombre} está aquí?",
                f"• Quiero unirme a {empresa} como {puesto} y generar impacto desde el primer día",
                "• Lo que busco: [tipo de impacto, reto o aprendizaje]",
                "",
                "## En 10 segundos:",
                f"• {anios} años de experiencia en [área especializada]  ·  especializado en [dominio]",
                "• Especializado en [habilidad o sector específico]",
                "• Apasionado por [tema relacionado con el rol]",
            ],
        }),
        ("timeline", {
            "titulo": "Mi Trayectoria",
            "steps": [
                (1, "[Inicio]\n[Empresa / Rol]", "[Año] — [Logro clave]"),
                (2, "[Crecimiento]\n[Empresa / Rol]", "[Año] — [Hito o responsabilidad]"),
                (3, "[Especialización]\n[Empresa / Rol]", "[Año] — [Proyecto destacado]"),
                (4, "[Actualidad]\n[Empresa / Rol]", "[Año] — [Situación presente]"),
                (5, "[Próximo paso]\n" + puesto, empresa),
            ],
        }),
        ("stats", {
            "titulo": "Mis Resultados — Lo que he logrado",
            "stats": [
                ("[X]%",  "Impacto logrado",       "Descripción del resultado más fuerte"),
                ("[X]",   "Proyectos liderados",   "En [N] años / en [empresa]"),
                ("[X]",   "Equipos / personas",    "Gestionados o liderados"),
                ("[X]+",  "Herramientas / Skills", "Stack técnico o metodologías"),
            ],
        }),
        ("two_col", {
            "titulo": "¿Por qué yo? — Mi propuesta de valor",
            "left_title": "LO QUE TRAIGO",
            "left_items": [
                "## Habilidades técnicas",
                "• [Skill 1] — nivel [básico/intermedio/experto]",
                "• [Skill 2] — nivel [básico/intermedio/experto]",
                "• [Skill 3] — nivel [básico/intermedio/experto]",
                "",
                "## Habilidades blandas",
                "• [Habilidad blanda 1]",
                "• [Habilidad blanda 2]",
            ],
            "right_title": f"LO QUE HARÉ EN {empresa.upper()[:20]}",
            "right_items": [
                "## En los primeros 30 días",
                "• [Acción de escucha / aprendizaje]",
                "",
                "## En los primeros 90 días",
                "• [Primera contribución concreta]",
                "• [Iniciativa o mejora que propondría]",
                "",
                "## A 1 año",
                "• [Resultado que esperas haber generado]",
            ],
        }),
        ("bullets", {
            "titulo": f"¿Por qué {empresa}?",
            "bullets": [
                f"## Por qué me atrae {empresa}",
                "",
                "• [Razón 1 — algo específico de la empresa: producto, misión, cultura]",
                "• [Razón 2 — cómo conecta con tu trayectoria o intereses]",
                "• [Razón 3 — algo que aprendiste investigando la empresa]",
                "",
                "## Preguntas que tengo",
                "• [Pregunta inteligente sobre el equipo, el reto o el rol]",
                "• [Pregunta sobre cultura o crecimiento]",
                "",
                f"## Resumen: soy {nombre}, traigo [X años + especialidad]",
                f"  y estoy listo para generar impacto en {empresa} como {puesto}.",
            ],
        }),
    ]


def _plan_academic(titulo, d, secciones, T, cfg=None):
    """Presentación académica — TFG, tesis, proyecto de clase, investigación.
    Estructura: intro → contexto → metodología → resultados → conclusiones."""
    autor    = d.get("autor") or d.get("speaker") or "[Autor/a]"
    tutor    = d.get("tutor") or ""
    institucion = d.get("empresa") or d.get("institucion") or "[Institución]"
    return [
        ("agenda", {
            "titulo_agenda": "ÍNDICE",
            "items": [
                (1, "Introducción y justificación", ""),
                (2, "Marco teórico / Antecedentes",  ""),
                (3, "Objetivos e hipótesis",          ""),
                (4, "Metodología",                   ""),
                (5, "Resultados",                    ""),
                (6, "Discusión y limitaciones",       ""),
                (7, "Conclusiones",                  ""),
                (8, "Referencias y agradecimientos",  ""),
            ],
        }),
        ("bullets", {
            "titulo": "Introducción y Justificación",
            "bullets": [
                f"## ¿Por qué este tema? — Motivación de {autor}",
                "• [Problema o fenómeno que motiva la investigación]",
                "• [Relevancia actual — por qué importa ahora]",
                "• [Gap en la literatura o práctica que este trabajo aborda]",
                "",
                "## Pregunta de investigación",
                "• [Pregunta central que guía todo el trabajo]",
                "",
                "## Alcance y delimitaciones",
                "• [Qué cubre y qué no cubre este trabajo]",
            ],
        }),
        ("two_col", {
            "titulo": "Marco Teórico y Antecedentes",
            "left_title": "CONCEPTOS CLAVE",
            "left_items": [
                "## Concepto 1",
                "• [Definición según autor, año]",
                "",
                "## Concepto 2",
                "• [Definición según autor, año]",
                "",
                "## Concepto 3",
                "• [Definición según autor, año]",
            ],
            "right_title": "ESTADO DEL ARTE",
            "right_items": [
                "## Investigaciones previas",
                "• [Estudio 1 — autor, año — hallazgo relevante]",
                "• [Estudio 2 — autor, año — hallazgo relevante]",
                "",
                "## Brechas identificadas",
                "• [Lo que la literatura aún no ha resuelto]",
                "• [Cómo este trabajo contribuye]",
            ],
        }),
        ("bullets", {
            "titulo": "Objetivos e Hipótesis",
            "bullets": [
                "## Objetivo general",
                "• [Objetivo principal del trabajo — verbo + qué + para qué]",
                "",
                "## Objetivos específicos",
                "• OE1: [objetivo específico 1]",
                "• OE2: [objetivo específico 2]",
                "• OE3: [objetivo específico 3]",
                "",
                "## Hipótesis / Preguntas derivadas",
                "• H1: [hipótesis o pregunta 1]",
                "• H2: [hipótesis o pregunta 2]",
            ],
        }),
        ("bullets", {
            "titulo": "Metodología",
            "bullets": [
                "## Enfoque de investigación",
                "• Tipo: [cuantitativo / cualitativo / mixto]",
                "• Diseño: [experimental / descriptivo / exploratorio / caso de estudio]",
                "",
                "## Muestra / Datos",
                "• Universo: [descripción de la población]",
                "• Muestra: [N] [unidades / participantes / registros]",
                "• Período: [fechas del estudio]",
                "",
                "## Técnicas e instrumentos",
                "• [Instrumento 1 — encuesta / entrevista / análisis documental]",
                "• [Instrumento 2 — software / herramienta de análisis]",
                "",
                "## Criterios de validez y confiabilidad",
                "• [Cómo se garantiza la calidad de los datos]",
            ],
        }),
        ("two_col", {
            "titulo": "Resultados",
            "left_title": "HALLAZGOS PRINCIPALES",
            "left_items": [
                "## Resultado 1",
                "• [Descripción + cifra o evidencia]",
                "• Relevancia: [por qué importa]",
                "",
                "## Resultado 2",
                "• [Descripción + cifra o evidencia]",
                "",
                "## Resultado 3",
                "• [Descripción + cifra o evidencia]",
            ],
            "right_title": "CONTRASTACIÓN DE HIPÓTESIS",
            "right_items": [
                "## H1",
                "• [Confirmada / Rechazada / Parcialmente confirmada]",
                "• [Explicación breve]",
                "",
                "## H2",
                "• [Confirmada / Rechazada]",
                "• [Explicación breve]",
                "",
                "## Hallazgos inesperados",
                "• [Algo que no esperabas encontrar]",
            ],
        }),
        ("bullets", {
            "titulo": "Conclusiones",
            "bullets": [
                "## Respuesta a la pregunta de investigación",
                "• [Síntesis de los hallazgos más importantes en 2-3 líneas]",
                "",
                "## Contribuciones del trabajo",
                "• Teórica: [aporte al conocimiento]",
                "• Práctica: [aplicación o recomendación para profesionales]",
                "",
                "## Limitaciones",
                "• [Limitación 1 — cómo afecta la generalización]",
                "• [Limitación 2 — sesgo o restricción metodológica]",
                "",
                "## Líneas futuras de investigación",
                "• [Pregunta que queda abierta para futuros trabajos]",
                "",
                f"Autor/a: {autor}" + (f"  ·  Tutor/a: {tutor}" if tutor else "") + f"  ·  {institucion}",
            ],
        }),
    ]


def _plan_client_onboarding(titulo, d, secciones, T, cfg=None):
    # ── v15.14.0 — Localización por país (moneda, autoridad fiscal) ──
    cfg = cfg or {}
    try:
        from agent.fiscal_verifier import verify_fiscal_data
        cfg = verify_fiscal_data(cfg)
    except Exception:
        pass
    sym = cfg.get("simbolo", "$")
    moneda = cfg.get("moneda", "USD")
    pais_nombre = cfg.get("pais", "")
    tax_authority = cfg.get("autoridad_fiscal", "")
    """Onboarding de cliente — primer deck que recibe un cliente al iniciar el servicio.
    Muy frecuente en agencias, consultoras y freelancers."""
    empresa  = d.get("empresa")  or "[Tu empresa / nombre]"
    cliente  = d.get("cliente")  or "[Nombre del cliente]"
    proyecto = d.get("proyecto") or d.get("producto") or "[Nombre del proyecto]"
    pm       = d.get("autor")    or d.get("pm") or "[Nombre del PM / responsable]"
    inicio   = d.get("inicio")   or "[Fecha de inicio]"
    return [
        ("bullets", {
            "titulo": f"Bienvenido/a, {cliente} 👋",
            "bullets": [
                f"## {empresa} × {cliente} — {proyecto}",
                "",
                f"Esta es la guía de referencia del proyecto {proyecto} entre {empresa} y {cliente}",
                f"## {proyecto}",
                "",
                f"• Cómo trabaja {empresa} y qué puede esperar {cliente}",
                "• Tu equipo dedicado y puntos de contacto",
                "• Cronograma, hitos y entregables",
                "• Lo que necesitamos de tu parte para avanzar",
                "",
                f"🚀 {proyecto}  ·  📅 {inicio}  ·  👤 PM: {pm}",
            ],
        }),
        ("two_col", {
            "titulo": "Tu Equipo Dedicado",
            "left_title": "NUESTRO LADO",
            "left_items": [
                f"## {empresa}",
                "",
                f"## {pm}",
                "Project Manager — tu punto de contacto principal",
                "",
                "## [Nombre]",
                "[Rol] — responsable de [área específica]",
                "",
                "## [Nombre]",
                "[Rol] — responsable de [área específica]",
            ],
            "right_title": "LO QUE NECESITAMOS DE TI",
            "right_items": [
                "## Tu equipo de contacto",
                "• [Nombre] — [Rol] — decisor principal",
                "• [Nombre] — [Rol] — contacto operativo",
                "",
                "## Accesos y recursos que necesitamos",
                "• [Acceso 1 — ej: Google Analytics, CRM]",
                "• [Acceso 2 — ej: servidor, repositorio]",
                "• [Documento o información — ej: brand guide]",
            ],
        }),
        ("timeline", {
            "titulo": f"Cronograma — {proyecto}",
            "steps": [
                (1, "Kick-of",        f"{inicio} — Alineación y setup"),
                (2, "Diagnóstico",     "Semana [N] — Análisis y definición"),
                (3, "Desarrollo",      "Semanas [N-N] — Ejecución principal"),
                (4, "Revisión",        "Semana [N] — Feedback y ajustes"),
                (5, "Entrega final",   "Fecha [X] — Cierre y handof"),
            ],
        }),
        ("bullets", {
            "titulo": "Cómo Trabajamos",
            "bullets": [
                "## Comunicación",
                "• Canal principal: [Slack / Email / WhatsApp / Teams]",
                "• Respuesta garantizada en: [N] horas hábiles",
                "• Reuniones de seguimiento: [frecuencia — semanal/quincenal]",
                "• Formato de reunión: [Zoom / presencial / híbrido]",
                "",
                f"## Gestión del proyecto {proyecto}",
                "• Herramienta: [Notion / Asana / Trello / ClickUp]",
                "• Compartiremos: actualizaciones de estado en tiempo real",
                "",
                "## Entregas y revisiones",
                "• Cada entrega viene con [N] ronda/s de revisión incluidas",
                "• Proceso: entregamos → feedback en [N] días → ajustamos",
                "• Cambios fuera de alcance: se gestionan vía Change Order",
            ],
        }),
        ("two_col", {
            "titulo": "Entregables y Criterios de Éxito",
            "left_title": "QUÉ VAS A RECIBIR",
            "left_items": [
                f"## Entregables de {proyecto}",
                "• [Entregable 1] — [Fecha]",
                "• [Entregable 2] — [Fecha]",
                "• [Entregable 3] — [Fecha]",
                "• [Entregable 4 — final] — [Fecha]",
                "",
                "## Formato de entrega",
                "• [Formato — PDF, drive, staging, etc.]",
            ],
            "right_title": "CÓMO SABREMOS QUE TUVIMOS ÉXITO",
            "right_items": [
                "## KPIs acordados",
                "• KPI 1: [métrica y target]",
                "• KPI 2: [métrica y target]",
                "• KPI 3: [métrica y target]",
                "",
                "## Criterios de aceptación",
                "• [Condición 1 para dar por cerrado el proyecto]",
                "• [Condición 2]",
            ],
        }),
        ("bullets", {
            "titulo": "Preguntas Frecuentes — FAQ",
            "bullets": [
                "## ¿Qué pasa si necesito cambiar algo del alcance?",
                "• Conversemos, siempre hay forma de ajustar. Los cambios menores",
                "  los absorbemos; los cambios de alcance se gestionan por Change Order.",
                "",
                "## ¿Cómo doy feedback a las entregas?",
                f"• Por {pm} vía [canal] — consolidamos el feedback antes de aplicar cambios.",
                "",
                f"## ¿Qué pasa si {proyecto} tiene un retraso?",
                "• Te avisamos con [N] días de anticipación — nunca te sorprendemos.",
                "",
                "## ¿Cómo escalo una preocupación?",
                f"• Directo con {pm} — y si es urgente, [canal de escalamiento].",
                "",
                "## ¿Cuándo arrancamos?",
                f"• {inicio} — {empresa} está listo para comenzar",
            ],
        }),
    ]


def _plan_all_hands(titulo, d, secciones, T, cfg=None):
    # ── v15.14.0 — Localización por país (moneda, autoridad fiscal) ──
    cfg = cfg or {}
    try:
        from agent.fiscal_verifier import verify_fiscal_data
        cfg = verify_fiscal_data(cfg)
    except Exception:
        pass
    sym = cfg.get("simbolo", "$")
    moneda = cfg.get("moneda", "USD")
    pais_nombre = cfg.get("pais", "")
    tax_authority = cfg.get("autoridad_fiscal", "")
    # ── v15.15.0 — LaborVerifier: actualización autónoma de prestaciones laborales ──
    try:
        from agent.fiscal_verifier import verify_labor_data
        cfg = verify_labor_data(cfg)
    except Exception:
        pass
    """All-hands / Town Hall — la dirección presenta al equipo completo.
    Formato motivacional con datos: visión, resultados, reconocimientos, retos."""
    empresa  = d.get("empresa")  or "[Empresa]"
    ceo      = d.get("autor")    or d.get("speaker") or "[CEO / Director]"
    periodo  = d.get("periodo")  or "[Período]"
    return [
        ("bullets", {
            "titulo": f"All-Hands {empresa} — {periodo}",
            "bullets": [
                "## Agenda de hoy",
                "",
                f"• 📊 Resultados {empresa} {periodo} — cómo nos fue",
                "• 🎯 Estado de nuestra estrategia — dónde estamos",
                "• 🌟 Reconocimientos — celebramos los logros del equipo",
                "• 🔭 Lo que viene — retos y oportunidades",
                "• 🗣 Preguntas y conversación abierta",
                "",
                f"All-Hands de {empresa}  ·  Presentado por {ceo}  ·  {periodo}",
            ],
        }),
        ("quote", {
            "quote":  "[Una frase que capture el espíritu del período o el mensaje central que quieres dejar hoy]",
            "author": ceo,
            "role":   f"CEO / Director General de {empresa}  ·  {periodo}",
        }),
        ("stats", {
            "titulo": f"Así nos fue — {periodo}",
            "stats": [
                ("[X]%",  "Resultado principal",   "vs objetivo de [X]%"),
                ("[X]",   "Clientes / Proyectos",  "Activos este período"),
                ("[X]%",  "NPS / Satisfacción",   "Encuesta de clientes"),
                ("[X]",   "Personas en el equipo", "y seguimos creciendo"),
            ],
        }),
        ("two_col", {
            "titulo": "Nuestra Estrategia — ¿Cómo vamos?",
            "left_title": "✅ LO QUE LOGRAMOS",
            "left_items": [
                "## Iniciativas completadas",
                "• [Iniciativa 1] — resultado concreto",
                "• [Iniciativa 2] — resultado concreto",
                "• [Iniciativa 3] — resultado concreto",
                "",
                "## Hitos del equipo",
                "• [Hito más relevante del período]",
                "• [Segundo hito importante]",
            ],
            "right_title": "🔄 EN PROGRESO Y APRENDIZAJES",
            "right_items": [
                "## Lo que está en marcha",
                "• [Iniciativa en curso 1]",
                "• [Iniciativa en curso 2]",
                "",
                "## Lo que aprendimos",
                "• [Aprendizaje honesto 1 — sin maquillaje]",
                "• [Aprendizaje honesto 2]",
                "",
                "## Ajustes que haremos",
                "• [Cambio de rumbo o prioridad]",
            ],
        }),
        ("bullets", {
            "titulo": "🌟 Reconocimientos — Gracias al Equipo",
            "bullets": [
                "## Personas que hicieron la diferencia este período",
                "",
                "## 🏆 [Nombre del equipo o persona]",
                "• [Descripción del logro — específico, no genérico]",
                "• Impacto: [cómo benefició al equipo, al cliente o a la empresa]",
                "",
                "## 🏆 [Nombre del equipo o persona]",
                "• [Descripción del logro]",
                "• Impacto: [descripción]",
                "",
                "## 🏆 [Nombre del equipo o persona]",
                "• [Descripción del logro]",
                "",
                f"## Todo el equipo de {empresa} — gracias por su compromiso diario.",
            ],
        }),
        ("timeline", {
            "titulo": "Lo que Viene — Próximos 90 Días",
            "steps": [
                (1, "Mes 1",    "[Prioridad principal — qué nos enfocamos]"),
                (2, "Mes 2",    "[Segunda prioridad — qué lanzamos o completamos]"),
                (3, "Mes 3",    "[Cierre del trimestre — meta a alcanzar]"),
                (4, "Hito",     "[Evento o lanzamiento importante próximo]"),
                (5, "Revisión", "[Próximo All-Hands — fecha]"),
            ],
        }),
        ("bullets", {
            "titulo": "Nuestra Cultura — Lo que Nos Define",
            "bullets": [
                "## Quiénes somos como equipo",
                "",
                "• [Valor 1 de la cultura — cómo se vive en el día a día]",
                "• [Valor 2 — ejemplo concreto reciente de cómo lo vivimos]",
                "• [Valor 3 — algo que queremos fortalecer]",
                "",
                "## Lo que nos hace únicos",
                "• [Característica del equipo que los distingue]",
                "",
                "## El compromiso de la dirección",
                f"• {ceo}: [Compromiso concreto con el equipo para el siguiente período]",
                "• [Recurso o mejora que la empresa se compromete a dar]",
            ],
        }),
        ("bullets", {
            "titulo": "Espacio Abierto — Tus Preguntas",
            "bullets": [
                "## Este es tu momento",
                "",
                "• No hay preguntas tontas",
                "• Si no podemos responder hoy, respondemos por escrito antes de [fecha]",
                "• Las preguntas anónimas también son bienvenidas: [canal]",
                "",
                "## Temas que podemos tocar:",
                "• Estrategia y dirección de la empresa",
                "• Cultura, equipo y bienestar",
                "• Procesos y herramientas",
                "• Lo que quieras que la dirección sepa",
                "",
                "## Encuesta de pulso post all-hands:",
                "• [Link] — 2 minutos, anónima, nos ayuda a mejorar",
            ],
        }),
    ]


def _plan_upsell(titulo, d, secciones, T, cfg=None):
    # ── v15.14.0 — Localización por país (moneda, autoridad fiscal) ──
    cfg = cfg or {}
    try:
        from agent.fiscal_verifier import verify_fiscal_data
        cfg = verify_fiscal_data(cfg)
    except Exception:
        pass
    sym = cfg.get("simbolo", "$")
    moneda = cfg.get("moneda", "USD")
    pais_nombre = cfg.get("pais", "")
    tax_authority = cfg.get("autoridad_fiscal", "")
    """Deck de upsell / renovación — para clientes actuales.
    Mix de resultados logrados + propuesta de expansión."""
    empresa  = d.get("empresa")  or "[Tu empresa]"
    cliente  = d.get("cliente")  or "[Nombre del cliente]"
    producto = d.get("producto") or "[Servicio actual]"
    nuevo    = d.get("nuevo")    or "[Nuevo servicio / expansión propuesta]"
    valor    = (d.get("monto") if d.get("monto") and "[" not in d.get("monto","") else None) or d.get("val") or "[Nuevo valor de la propuesta]"
    try:
        cbuf = _chart_bars(
            ["Resultado 1","Resultado 2","Resultado 3","Resultado 4"],
            [72, 88, 95, 81],
            f"Resultados logrados — {cliente} (%)",
            "#0F6E56", "F0F4F8"
        )
    except Exception:
        cbuf = None
    slides = [
        ("bullets", {
            "titulo": f"Un año juntos — Lo que logramos, {cliente}",
            "bullets": [
                f"## {empresa} × {cliente} — Lo que logramos juntos",
                f"Resumen de resultados del período trabajando juntos",
                "",
                f"• Resultado principal de {empresa} para {cliente}: [cifra concreta]",
                "• [Segundo resultado — con contexto]",
                "• [Tercer resultado — tendencia positiva]",
                "",
                "## Lo que esto significa para tu negocio",
                "• [Impacto en revenue / costos / eficiencia]",
                "• [Impacto en equipo o clientes del cliente]",
            ],
        }),
    ]
    if cbuf:
        slides.append(("chart", {
            "titulo": f"Resultados en Números — {producto}",
            "chart_buf": cbuf,
            "caption": "Indicadores clave del período de colaboración",
            "bullets": [
                "## Análisis",
                "• Mejor resultado: [área] con [X]%",
                "• Tendencia: consistentemente al alza",
                "• vs objetivo inicial: superado en [X]%",
                "",
                "## Qué lo hizo posible",
                "• [Factor 1]",
                "• [Factor 2]",
            ],
        }))
    slides += [
        ("quote", {
            "quote":  f"[Cita real o parafraseo de feedback positivo que {cliente} te ha dado sobre el trabajo]",
            "author": "[Nombre del contacto principal]",
            "role":   f"[Cargo] — {cliente}",
        }),
        ("two_col", {
            "titulo": "La Oportunidad que Vemos",
            "left_title": "DÓNDE ESTÁS HOY",
            "left_items": [
                "## Lo que ya tienes funcionando",
                f"• {producto}: generando [resultado]",
                "• Equipo alineado con el proceso",
                "• Métricas en verde",
                "",
                "## El techo que estamos tocando",
                "• [Limitación del alcance actual]",
                "• [Oportunidad que no estamos capturando aún]",
            ],
            "right_title": "A DÓNDE PODEMOS LLEVARTE",
            "right_items": [
                "## La oportunidad",
                f"• {nuevo}",
                "",
                "## Por qué ahora",
                "• [Razón de timing — mercado, madurez, momento interno]",
                "",
                "## Resultado esperado",
                "• [Proyección concreta con cifra]",
                "• [Plazo estimado para ver resultados]",
            ],
        }),
        ("bullets", {
            "titulo": f"La Propuesta — {nuevo}",
            "bullets": [
                f"## ¿Qué incluye {nuevo}?",
                "",
                "• [Componente 1 — descripción y beneficio]",
                "• [Componente 2 — descripción y beneficio]",
                "• [Componente 3 — descripción y beneficio]",
                "",
                "## Condiciones especiales para clientes actuales",
                f"• Inversión: {valor}  (condiciones especiales para clientes actuales)",
                "• [Descuento o beneficio por ser cliente existente]",
                "• [Flexibilidad de pago o condiciones especiales]",
                "",
                "## Timeline de implementación",
                "• Inicio: [fecha]  ·  Primeros resultados: [N] semanas",
            ],
        }),
        ("bullets", {
            "titulo": "Próximos Pasos",
            "bullets": [
                "## Para avanzar juntos",
                "",
                f"• Paso 1 — Esta semana: ¿Alguna pregunta sobre la propuesta?",
                f"• Paso 2 — [Fecha]: Call de alineación (30 min)",
                f"• Paso 3 — [Fecha]: Firma del addendum / nuevo contrato",
                f"• Paso 4 — [Fecha]: Inicio de {nuevo}",
                "",
                "## Por qué tiene sentido seguir juntos",
                f"• Ya conocemos el negocio de {cliente}, su equipo y sus procesos",
                "• Cero curva de aprendizaje — arrancamos desde donde estamos",
                "• El historial habla por sí solo",
                "",
                f"¿Avanzamos juntos? — {empresa} × {cliente}",
            ],
        }),
    ]
    return slides

# ─────────────────────────────────────────────────────────────────────────────
# TEMPLATE REGISTRY (v14.2) — al final del archivo para evitar forward refs
# ─────────────────────────────────────────────────────────────────────────────
# Mapeo explícito template_id → planner. El dispatcher de generate_powerpoint
# lo consulta primero; si no hay match, cae al matching por keywords (legacy).
PPT_TEMPLATES = {
    "pitch":              _plan_pitch,
    "sales":              _plan_sales,
    "report":             _plan_report,
    "training":           _plan_training,
    "company":            _plan_company,
    "product":            _plan_product,
    "strategy":           _plan_strategy,
    "board":              _plan_board,
    "investor_update":    _plan_investor_update,
    "case_study":         _plan_case_study,
    "webinar":            _plan_webinar,
    "personal_brand":     _plan_personal_brand,
    "job_presentation":   _plan_job_presentation,
    "academic":           _plan_academic,
    "client_onboarding":  _plan_client_onboarding,
    "all_hands":          _plan_all_hands,
    "upsell":             _plan_upsell,
    "generic":            _plan_generic,
}
